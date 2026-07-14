"""典型文件示例生成器。

为各支持格式动态生成代表性测试文件，每个文件注入 ``password`` 关键词。
用于功能测试（验证提取正确性）和性能基准（测量提取速度）。

可生成格式：
    - 纯文本：txt/json/yaml/xml/csv/md/html
    - 二进制：rtf/docx/xlsx/pptx/eml

不可生成格式（需外部工具或无法可靠生成）：
    pdf/doc/ppt/xls/msg/odt/ods
"""

from __future__ import annotations

import random
from email.message import EmailMessage
from pathlib import Path
from typing import Callable

__all__ = [
    "BINARY_GENERATORS",
    "GENERATORS",
    "TEXT_GENERATORS",
    "generate_file",
    "generate_files",
    "generate_sample_bytes",
]

# 敏感数据样本（每个文件注入其一，均含 password 关键词）
_SECRETS: tuple[str, ...] = (
    "password=secret123",
    "user_password=admin",
    "db_password=rootpass",
)

# 填充文本模板
_FILLER = "the quick brown fox jumps over the lazy dog\n"

# 纯文本格式 → 生成函数映射
TEXT_GENERATORS: dict[str, Callable[[Path, str], None]] = {}

# 二进制格式 → 生成函数映射
BINARY_GENERATORS: dict[str, Callable[[Path, str], None]] = {}

# 全部可生成格式
GENERATORS: dict[str, Callable[[Path, str], None]] = {}


def _register(ext: str, binary: bool = False) -> Callable[[Callable[[Path, str], None]], Callable[[Path, str], None]]:
    """注册格式生成器的装饰器。"""

    def decorator(fn: Callable[[Path, str], None]) -> Callable[[Path, str], None]:
        target = BINARY_GENERATORS if binary else TEXT_GENERATORS
        target[ext] = fn
        GENERATORS[ext] = fn
        return fn

    return decorator


def _make_filler_text(size: int, rng: random.Random) -> str:
    """生成指定大小的填充文本，首行注入敏感数据。"""
    lines: list[str] = []
    written = 0
    secret = rng.choice(_SECRETS)
    lines.append(f"# {secret}\n")
    written += len(secret) + 3
    while written < size:
        lines.append(_FILLER)
        written += len(_FILLER)
    return "".join(lines)[:size]


# ---------------------------------------------------------------------------
# 纯文本格式
# ---------------------------------------------------------------------------


@_register("txt")
def _gen_txt(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


@_register("json")
def _gen_json(path: Path, content: str) -> None:
    """生成包含文本内容的 JSON 文件。"""
    import json

    lines = [line for line in content.split("\n") if line]
    data = {"content": "\n".join(lines), "items": lines[:10]}
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


@_register("yaml")
def _gen_yaml(path: Path, content: str) -> None:
    """生成包含文本内容的 YAML 文件。"""
    lines = [line for line in content.split("\n") if line]
    yaml_text = "content: |\n"
    for line in lines:
        yaml_text += f"  {line}\n"
    yaml_text += "items:\n"
    for line in lines[:10]:
        yaml_text += f"  - {line}\n"
    path.write_text(yaml_text, encoding="utf-8")


@_register("xml")
def _gen_xml(path: Path, content: str) -> None:
    """生成包含文本内容的 XML 文件。"""
    lines = content.split("\n")
    xml = '<?xml version="1.0" encoding="UTF-8"?>\n<document>\n'
    for i, line in enumerate(lines):
        if line:
            xml += f'  <line id="{i}">{line}</line>\n'
    xml += "</document>\n"
    path.write_text(xml, encoding="utf-8")


@_register("csv")
def _gen_csv(path: Path, content: str) -> None:
    """生成包含文本内容的 CSV 文件。"""
    lines = [line for line in content.split("\n") if line]
    csv_text = "id,content\n"
    for i, line in enumerate(lines):
        csv_text += f'{i},"{line}"\n'
    path.write_text(csv_text, encoding="utf-8")


@_register("md")
def _gen_md(path: Path, content: str) -> None:
    """生成包含文本内容的 Markdown 文件。"""
    lines = [line for line in content.split("\n") if line]
    md = "# 文档标题\n\n"
    for line in lines:
        md += f"- {line}\n"
    path.write_text(md, encoding="utf-8")


@_register("html")
def _gen_html(path: Path, content: str) -> None:
    """生成包含文本内容的 HTML 文件。"""
    lines = [line for line in content.split("\n") if line]
    html = "<html><body>\n"
    for line in lines:
        html += f"<p>{line}</p>\n"
    html += "</body></html>\n"
    path.write_text(html, encoding="utf-8")


# ---------------------------------------------------------------------------
# 二进制格式
# ---------------------------------------------------------------------------


@_register("rtf", binary=True)
def _gen_rtf(path: Path, content: str) -> None:
    """生成包含文本内容的 RTF 文件。"""
    lines = [line for line in content.split("\n") if line]
    rtf = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}\f0\fs24 "
    for line in lines:
        rtf += line.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}") + r"\par "
    rtf += "}"
    path.write_text(rtf, encoding="utf-8")


@_register("docx", binary=True)
def _gen_docx(path: Path, content: str) -> None:
    """生成包含文本内容的 DOCX 文件。"""
    from docx import Document

    doc = Document()
    for line in content.split("\n"):
        if line:
            doc.add_paragraph(line)
    doc.save(str(path))


@_register("xlsx", binary=True)
def _gen_xlsx(path: Path, content: str) -> None:
    """生成包含文本内容的 XLSX 文件。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "数据"
    ws["A1"] = "内容"
    for i, line in enumerate(content.split("\n"), start=2):
        if line:
            ws[f"A{i}"] = line
    wb.save(str(path))


@_register("pptx", binary=True)
def _gen_pptx(path: Path, content: str) -> None:
    """生成包含文本内容的 PPTX 文件。"""
    from pptx import Presentation

    prs = Presentation()
    lines = [line for line in content.split("\n") if line]
    # 每页放 5 行
    chunk_size = 5
    for i in range(0, len(lines), chunk_size):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        chunk = lines[i : i + chunk_size]
        if slide.shapes.title:
            slide.shapes.title.text = chunk[0]
        if len(chunk) > 1 and len(slide.placeholders) > 1:
            slide.placeholders[1].text = "\n".join(chunk[1:])
    prs.save(str(path))


@_register("eml", binary=True)
def _gen_eml(path: Path, content: str) -> None:
    """生成包含文本内容的 EML 邮件文件。"""
    lines = [line for line in content.split("\n") if line]
    msg = EmailMessage()
    msg["Subject"] = "测试邮件"
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    msg.set_content("\n".join(lines))
    path.write_bytes(msg.as_bytes())


# ---------------------------------------------------------------------------
# 公共 API
# ---------------------------------------------------------------------------


def generate_file(path: Path, ext: str, size_hint: int = 4096, rng: random.Random | None = None) -> Path:
    """生成单个指定格式的示例文件。

    :param path: 输出文件路径
    :param ext: 扩展名（不含点，如 ``"docx"``）
    :param size_hint: 内容大小提示（字节），实际文件大小因格式而异
    :param rng: 随机数生成器；None 时新建
    :return: 生成的文件路径
    :raises ValueError: 不支持的格式
    """
    if ext not in GENERATORS:
        raise ValueError(f"不支持的格式: {ext}，可生成格式: {sorted(GENERATORS)}")
    if rng is None:
        rng = random.Random(42)
    content = _make_filler_text(size_hint, rng)
    GENERATORS[ext](path, content)
    return path


def generate_files(root: Path, count: int, seed: int = 42) -> list[Path]:
    """生成混合格式测试文件集，约 30% 含敏感数据。

    纯文本与二进制格式均匀分布，每个文件均注入 ``password`` 关键词。

    :param root: 输出目录（自动创建）
    :param count: 文件数量
    :param seed: 随机种子（可复现）
    :return: 生成文件路径列表
    """
    root.mkdir(parents=True, exist_ok=True)
    rng = random.Random(seed)
    all_exts = list(GENERATORS.keys())
    paths: list[Path] = []
    for i in range(count):
        ext = rng.choice(all_exts)
        path = root / f"file_{i:05d}.{ext}"
        size = rng.randint(1024, 30 * 1024)
        generate_file(path, ext, size, rng)
        paths.append(path)
    return paths


def generate_sample_bytes(ext: str, size_hint: int = 4096, rng: random.Random | None = None) -> bytes:
    """生成指定格式的示例文件字节（内存中），用于提取器速度测试。

    :param ext: 扩展名（不含点）
    :param size_hint: 内容大小提示
    :param rng: 随机数生成器；None 时新建
    :return: 文件字节内容
    :raises ValueError: 不支持的格式
    """
    import tempfile

    if rng is None:
        rng = random.Random(42)
    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        generate_file(tmp_path, ext, size_hint, rng)
        return tmp_path.read_bytes()
    finally:
        tmp_path.unlink(missing_ok=True)
