"""Microsoft Office 文档提取器：DOCX 与 PPTX。

DOCX 使用 python-docx 提取段落、表格、页眉页脚。
PPTX 使用 python-pptx 提取幻灯片文本框、表格、备注。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List, Tuple

from pyfilescan.extractors.base import Extractor, ExtractorError

__all__ = ["DocxExtractor", "PptxExtractor"]

logger = logging.getLogger(__name__)


class DocxExtractor(Extractor):
    """DOCX 文档文本提取器。"""

    @property
    def supported_extensions(self) -> Tuple[str, ...]:
        return ("docx",)

    def extract(self, path: Path) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise ExtractorError("python-docx 未安装，无法提取 DOCX") from exc

        try:
            doc = Document(str(path))
        except Exception as exc:
            raise ExtractorError(f"DOCX 解析失败: {path}: {exc}") from exc

        parts: List[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append("\t".join(row_texts))

        for section in doc.sections:
            for header_footer in (section.header, section.footer):
                for para in header_footer.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        return "\n".join(parts)


class PptxExtractor(Extractor):
    """PPTX 演示文稿文本提取器。"""

    @property
    def supported_extensions(self) -> Tuple[str, ...]:
        return ("pptx",)

    def extract(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError("python-pptx 未安装，无法提取 PPTX") from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ExtractorError(f"PPTX 解析失败: {path}: {exc}") from exc

        parts: List[str] = []
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_texts = self._extract_slide(slide)
            if slide_texts:
                parts.append(f"--- 幻灯片 {slide_index} ---")
                parts.extend(slide_texts)

        return "\n".join(parts)

    def _extract_slide(self, slide: object) -> List[str]:
        """提取单张幻灯片的文本。"""
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append("\t".join(row_texts))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                texts.append(f"[备注] {notes_text}")
        return texts
