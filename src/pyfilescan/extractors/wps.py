"""WPS Office 文档提取器。

WPS 文档格式（.wps/.et/.dps）有两种形态：

1. **OOXML 兼容**：现代 WPS 默认保存格式，本质是 ZIP 打包的 XML，
   可复用 python-docx/openpyxl/python-pptx 提取。
2. **旧版二进制**：早期 WPS 私有格式，无法用 Office 库解析，
   记录告警并返回空字符串。

通过检查文件头是否为 ZIP 魔数（PK\\x03\\x04）判断格式类型。
"""

from __future__ import annotations

import logging
import os
import shutil
import tempfile
from pathlib import Path
from typing import Tuple

from pyfilescan.extractors.base import Extractor, ExtractorError

__all__ = ["WpsExtractor"]

logger = logging.getLogger(__name__)

_ZIP_MAGIC = b"PK\x03\x04"


class WpsExtractor(Extractor):
    """WPS Office 文档提取器，按扩展名分发到对应的 OOXML 提取逻辑。"""

    @property
    def supported_extensions(self) -> Tuple[str, ...]:
        return ("wps", "et", "dps")

    def extract(self, path: Path) -> str:
        if not self._is_ooxml(path):
            logger.info("WPS 文档非 OOXML 格式（旧版二进制），跳过: %s", path)
            return ""

        ext = path.suffix.lower().lstrip(".")
        if ext == "wps":
            return self._extract_as_docx(path)
        if ext == "et":
            return self._extract_as_xlsx(path)
        if ext == "dps":
            return self._extract_as_pptx(path)
        return ""

    def _is_ooxml(self, path: Path) -> bool:
        """检查文件头是否为 ZIP 魔数（OOXML 是 ZIP 打包）。"""
        try:
            with path.open("rb") as fh:
                return fh.read(4) == _ZIP_MAGIC
        except OSError:
            return False

    def _temp_with_ext(self, path: Path, target_ext: str) -> Path:
        """复制文件到临时路径并改用目标扩展名，绕过库的扩展名检查。"""
        fd, tmp_name = tempfile.mkstemp(suffix=f".{target_ext}")
        tmp_path = Path(tmp_name)
        # mkstemp 返回的 fd 必须先关闭，否则 Windows 上文件被锁定
        os.close(fd)
        try:
            shutil.copyfile(path, tmp_path)
        except OSError:
            self._safe_unlink(tmp_path)
            raise
        return tmp_path

    @staticmethod
    def _safe_unlink(path: Path) -> None:
        """安全删除文件，忽略 Windows 上的文件锁定错误。"""
        try:
            path.unlink(missing_ok=True)
        except PermissionError:
            logger.debug("临时文件被锁定，跳过删除: %s", path)

    def _extract_as_docx(self, path: Path) -> str:
        """以 DOCX 方式提取 WPS 文字文档。"""
        try:
            from docx import Document
        except ImportError as exc:
            raise ExtractorError("python-docx 未安装，无法提取 WPS 文字文档") from exc

        tmp = self._temp_with_ext(path, "docx")
        try:
            doc = Document(str(tmp))
        except Exception as exc:
            raise ExtractorError(f"WPS 文字文档解析失败: {path}: {exc}") from exc
        finally:
            self._safe_unlink(tmp)

        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append("\t".join(row_texts))
        return "\n".join(parts)

    def _extract_as_xlsx(self, path: Path) -> str:
        """以 XLSX 方式提取 WPS 表格文档。"""
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ExtractorError("openpyxl 未安装，无法提取 WPS 表格") from exc

        tmp = self._temp_with_ext(path, "xlsx")
        try:
            wb = load_workbook(str(tmp), read_only=True, data_only=True)
        except Exception as exc:
            raise ExtractorError(f"WPS 表格解析失败: {path}: {exc}") from exc
        finally:
            self._safe_unlink(tmp)

        parts = []
        try:
            for sheet in wb:
                parts.append(f"--- 工作表: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    cell_texts = [str(c).strip() for c in row if c is not None]
                    if cell_texts:
                        parts.append("\t".join(cell_texts))
        finally:
            wb.close()
        return "\n".join(parts)

    def _extract_as_pptx(self, path: Path) -> str:
        """以 PPTX 方式提取 WPS 演示文档。"""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError("python-pptx 未安装，无法提取 WPS 演示") from exc

        tmp = self._temp_with_ext(path, "pptx")
        try:
            prs = Presentation(str(tmp))
        except Exception as exc:
            raise ExtractorError(f"WPS 演示解析失败: {path}: {exc}") from exc
        finally:
            self._safe_unlink(tmp)

        parts = []
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"--- 幻灯片 {slide_index} ---")
                parts.extend(slide_texts)
        return "\n".join(parts)
