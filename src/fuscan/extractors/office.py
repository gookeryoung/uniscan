"""Microsoft Office 文档提取器：DOCX 与 PPTX。

DOCX 使用 python-docx 提取段落、表格、页眉页脚。
PPTX 使用 python-pptx 提取幻灯片文本框、表格、备注。
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

from typing_extensions import override

from fuscan.extractors.base import Extractor, ExtractorError

__all__ = ["DocxExtractor", "PptxExtractor"]

logger = logging.getLogger(__name__)


class DocxExtractor(Extractor):
    """DOCX 文档文本提取器。"""

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        """返回 DOCX 提取器支持的扩展名。"""
        return ("docx",)

    @override
    def extract(self, path: Path) -> str:
        """提取 DOCX 段落、表格与页眉页脚文本。"""
        try:
            data = path.read_bytes()
        except OSError as exc:
            raise ExtractorError(f"文件读取失败: {path}: {exc}") from exc
        return self.extract_from_bytes(data)

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        """从内存字节提取 DOCX 文本。"""
        try:
            from docx import Document
        except ImportError as exc:
            raise ExtractorError("python-docx 未安装，无法提取 DOCX") from exc

        try:
            doc = Document(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"DOCX 解析失败: {exc}") from exc

        parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append("\t".join(row_texts))

        for section in doc.sections:
            for header_footer in (section.header, section.footer):
                for para in header_footer.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        return "\n".join(parts)


class PptxExtractor(Extractor):
    """PPTX 演示文稿文本提取器。"""

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        """返回 PPTX 提取器支持的扩展名。"""
        return ("pptx",)

    @override
    def extract(self, path: Path) -> str:
        """提取 PPTX 幻灯片文本框、表格与备注文本。"""
        try:
            data = path.read_bytes()
        except OSError as exc:
            raise ExtractorError(f"文件读取失败: {path}: {exc}") from exc
        return self.extract_from_bytes(data)

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        """从内存字节提取 PPTX 文本。"""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError("python-pptx 未安装，无法提取 PPTX") from exc

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"PPTX 解析失败: {exc}") from exc

        parts: list[str] = []
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_texts = self._extract_slide(slide)
            if slide_texts:
                parts.append(f"--- 幻灯片 {slide_index} ---")
                parts.extend(slide_texts)

        return "\n".join(parts)

    def _extract_slide(self, slide: object) -> list[str]:
        """提取单张幻灯片的文本。"""
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append("\t".join(row_texts))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                texts.append(f"[备注] {notes_text}")
        return texts
