"""WPS Office 文档提取器。

WPS 文档格式（.wps/.et/.dps）有两种形态：

1. **OOXML 兼容**：现代 WPS 默认保存格式，本质是 ZIP 打包的 XML，
   可复用 python-docx/openpyxl/python-pptx 提取。
2. **旧版二进制**：早期 WPS 私有格式，无法用 Office 库解析，
   记录告警并返回空字符串。

通过检查文件头是否为 ZIP 魔数（PK\\x03\\x04）判断格式类型，
再通过 ZIP 内部条目名区分 DOCX/XLSX/PPTX 子类型。
"""

from __future__ import annotations

import io
import logging
import zipfile
from pathlib import Path

from typing_extensions import override

from fuscan.extractors.base import Extractor, ExtractorError

__all__ = ["WpsExtractor"]

logger = logging.getLogger(__name__)

_ZIP_MAGIC = b"PK\x03\x04"


class WpsExtractor(Extractor):
    """WPS Office 文档提取器，按 OOXML 子类型分发到对应提取逻辑。"""

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        """返回 WPS 提取器支持的扩展名。"""
        return ("wps", "et", "dps")

    @override
    @property
    def display_name(self) -> str:
        """返回提取器的中文显示名称。"""
        return "WPS 文档"

    @override
    def extract(self, path: Path) -> str:
        """提取 WPS 文档文本，按 OOXML 子类型分发到对应提取逻辑。"""
        try:
            data = path.read_bytes()
        except OSError as exc:
            raise ExtractorError(f"文件读取失败: {path}: {exc}") from exc
        return self.extract_from_bytes(data)

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        """从内存字节提取 WPS 文档文本。"""
        if not data.startswith(_ZIP_MAGIC):
            logger.info("WPS 文档非 OOXML 格式（旧版二进制），跳过")
            return ""

        ooxml_type = self._detect_ooxml_type(data)
        if ooxml_type == "docx":
            return self._extract_as_docx(data)
        if ooxml_type == "xlsx":
            return self._extract_as_xlsx(data)
        if ooxml_type == "pptx":
            return self._extract_as_pptx(data)
        logger.info("WPS 文档 OOXML 子类型未知，跳过")
        return ""

    def _detect_ooxml_type(self, data: bytes) -> str | None:
        """通过 ZIP 内部条目名判断 OOXML 子类型。

        :return: ``"docx"`` / ``"xlsx"`` / ``"pptx"``，无法判断时返回 ``None``
        """
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = set(zf.namelist())
        except zipfile.BadZipFile:
            return None
        if "word/document.xml" in names:
            return "docx"
        if "xl/workbook.xml" in names:
            return "xlsx"
        if "ppt/presentation.xml" in names:
            return "pptx"
        return None

    def _extract_as_docx(self, data: bytes) -> str:
        """以 DOCX 方式提取 WPS 文字文档。"""
        try:
            from docx import Document
        except ImportError as exc:
            raise ExtractorError("python-docx 未安装，无法提取 WPS 文字文档") from exc

        try:
            doc = Document(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"WPS 文字文档解析失败: {exc}") from exc

        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append("\t".join(row_texts))
        return "\n".join(parts)

    def _extract_as_xlsx(self, data: bytes) -> str:
        """以 XLSX 方式提取 WPS 表格文档。"""
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ExtractorError("openpyxl 未安装，无法提取 WPS 表格") from exc

        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        except Exception as exc:
            raise ExtractorError(f"WPS 表格解析失败: {exc}") from exc

        parts = []
        try:
            for sheet in wb:
                parts.append(f"--- 工作表: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    cell_texts = [str(c).strip() for c in row if c is not None]
                    if cell_texts:
                        parts.append("\t".join(cell_texts))
        finally:
            wb.close()
        return "\n".join(parts)

    def _extract_as_pptx(self, data: bytes) -> str:
        """以 PPTX 方式提取 WPS 演示文档。"""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError("python-pptx 未安装，无法提取 WPS 演示") from exc

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"WPS 演示解析失败: {exc}") from exc

        parts = []
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:  # pyrefly: ignore [missing-attribute]
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"--- 幻灯片 {slide_index} ---")
                parts.extend(slide_texts)
        return "\n".join(parts)
