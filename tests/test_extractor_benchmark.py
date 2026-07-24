"""提取器解析速度基准测试（iter-90）。

为每类文件类型设计基准测试，验证 ``speed_tier`` 声明与实测性能一致。
所有测试标记 ``@pytest.mark.slow``，CI 默认跳过（``-m "not slow"``）。

5 档速度分级（``SpeedTier``）：

- T1 极速（``VERY_FAST``）：< 10ms/MB，纯字节解码
  - 纯文本/源代码/配置文件/标记与数据/样式表 5 个子提取器
- T2 快速（``FAST``）：10-50ms/MB，标准库解析
  - EML 邮件
- T3 中速（``MEDIUM``）：50-200ms/MB，单次 XML 解析 + 树遍历
  - DOCX/ODT/RTF/WPS/MSG
- T4 慢速（``SLOW``）：200-1000ms/MB，单元格遍历或字节级扫描
  - XLSX/ODS/XLS/PPTX/DOC/PPT
- T5 极慢（``VERY_SLOW``）：> 1000ms/MB，复杂页面布局分析
  - PDF

基准测试设计原则：

1. 生成 100KB-1MB 量级的典型样本文件（避免空文件或过小文件测量噪声）
2. 用 ``time.perf_counter`` 测量 3 次取中位数，降低抖动
3. 阈值宽松（为 CI 环境留 5-10 倍余量），仅验证档次声明合理
4. 打印测量结果供文档参考，不作为回归门禁
"""

from __future__ import annotations

import io
import logging
import statistics
import time
from collections.abc import Callable
from typing import Any

import pytest

from fuscan.extractors import (
    DocxExtractor,
    EmlExtractor,
    MsgExtractor,
    OdsExtractor,
    OdtExtractor,
    PdfExtractor,
    PlainTextExtractor,
    PptxExtractor,
    RtfExtractor,
    SourceCodeExtractor,
    WpsExtractor,
    XlsExtractor,
    XlsxExtractor,
)
from fuscan.extractors.base import Extractor, SpeedTier

logger = logging.getLogger(__name__)

# 测量次数：取中位数降低抖动
_ITERATIONS = 3

# 各档次耗时上限（秒），用于断言档次声明合理
# 阈值为典型样本（100KB-1MB）的宽松上限，CI 环境留 5-10 倍余量
_TIER_TIME_LIMITS: dict[SpeedTier, float] = {
    SpeedTier.VERY_FAST: 0.5,  # 100KB 文本 < 500ms
    SpeedTier.FAST: 1.0,  # 典型 EML < 1s
    SpeedTier.MEDIUM: 2.0,  # 典型 OOXML < 2s
    SpeedTier.SLOW: 5.0,  # 典型表格/幻灯片 < 5s
    SpeedTier.VERY_SLOW: 10.0,  # 典型 PDF < 10s
}


def _measure(extract_fn: Callable[..., Any], *args: Any, iterations: int = _ITERATIONS) -> float:
    """测量函数执行耗时（秒），返回中位数。

    :param extract_fn: 提取函数（如 ``extractor.extract_from_bytes``）
    :param args: 传给提取函数的参数
    :param iterations: 测量次数
    :return: 中位数耗时（秒）
    """
    times: list[float] = []
    for _ in range(iterations):
        start = time.perf_counter()
        extract_fn(*args)
        times.append(time.perf_counter() - start)
    return statistics.median(times)


def _assert_tier(extractor: Extractor, expected: SpeedTier) -> None:
    """断言提取器的 speed_tier 与预期档次一致。"""
    assert extractor.speed_tier == expected, (
        f"{type(extractor).__name__}.speed_tier 应为 {expected.label}，实际为 {extractor.speed_tier.label}"
    )


def _assert_time_within_tier(elapsed: float, tier: SpeedTier, extractor_name: str) -> None:
    """断言耗时在档次预期范围内（宽松阈值，避免 CI flakiness）。"""
    limit = _TIER_TIME_LIMITS[tier]
    assert elapsed < limit, (
        f"{extractor_name} 声明为 {tier.label}，但耗时 {elapsed:.3f}s 超过阈值 {limit}s。"
        f"如属环境差异，请放宽阈值或检查实现。"
    )


# ----------------------------- 样本生成工具 -----------------------------


def _make_text_sample(size_kb: int = 100) -> bytes:
    """生成指定大小的 UTF-8 文本样本（重复典型文本模式）。"""
    pattern = "hello password world 第二行内容 secret_key=abc123\n"
    repeat = max(1, size_kb * 1024 // len(pattern.encode("utf-8")) + 1)
    return (pattern * repeat).encode("utf-8")[: size_kb * 1024]


def _make_eml_sample() -> bytes:
    """生成典型 EML 邮件样本。"""
    import email.message
    import email.policy

    msg = email.message.EmailMessage(policy=email.policy.default)
    msg["Subject"] = "测试邮件 password secret"
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    body = "邮件正文含 password 与 secret_key\n" * 50
    msg.set_content(body)
    return msg.as_bytes()


def _make_docx_sample() -> bytes:
    """生成典型 DOCX 文档样本（含段落与表格）。"""
    from docx import Document

    doc = Document()
    for i in range(20):
        doc.add_paragraph(f"段落 {i}：含 password 和 secret 关键词")
    table = doc.add_table(rows=5, cols=3)
    for row in table.rows:
        for cell in row.cells:
            cell.text = "password cell"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_sample() -> bytes:
    """生成典型 XLSX 工作簿样本（多工作表、100 行 × 10 列）。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "数据"
    for row in range(100):
        for col in range(10):
            ws.cell(row=row + 1, column=col + 1, value=f"cell_{row}_{col}_password")
    wb.create_sheet("第二表")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_sample() -> bytes:
    """生成典型 PPTX 演示文稿样本（5 张幻灯片）。"""
    from pptx import Presentation

    prs = Presentation()
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"幻灯片 {i} 含 password"  # pyrefly: ignore [missing-attribute]
        slide.placeholders[1].text = f"内容 secret_key {i}"  # pyrefly: ignore [missing-attribute]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_rtf_sample() -> bytes:
    """生成典型 RTF 富文本样本。"""
    # 简单 RTF 格式：含多段落文本
    text = "password secret 内容" * 50
    rtf = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 SimSun;}} \f0\fs24 " + text + r"}"
    return rtf.encode("utf-8")


def _make_pdf_sample() -> bytes:
    """生成最小 PDF 样本（单页文本）。

    使用 reportlab 生成；未安装时跳过测试。
    """
    try:
        from reportlab.pdfgen import canvas
    except ImportError:
        pytest.skip("reportlab 未安装，跳过 PDF 基准测试")
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.setFont("Helvetica", 12)
    for i in range(20):
        c.drawString(100, 800 - i * 20, f"Line {i}: password secret content")
    c.save()
    return buf.getvalue()


def _make_odt_sample() -> bytes:
    """生成典型 ODT 文档样本。

    使用 odfpy 生成；未安装时跳过测试。
    """
    try:
        from odf.opendocument import OpenDocumentText
        from odf.text import P
    except ImportError:
        pytest.skip("odfpy 未安装，跳过 ODT 基准测试")
    doc = OpenDocumentText()
    for i in range(20):
        p = P(text=f"段落 {i}：password secret 内容")
        doc.text.addElement(p)  # type: ignore[attr-defined]
    buf = io.BytesIO()
    doc.write(buf)
    return buf.getvalue()


def _make_ods_sample() -> bytes:
    """生成典型 ODS 表格样本。

    使用 odfpy 生成；未安装时跳过测试。
    """
    try:
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableCell, TableRow
        from odf.text import P
    except ImportError:
        pytest.skip("odfpy 未安装，跳过 ODS 基准测试")
    doc = OpenDocumentSpreadsheet()
    table = Table(name="数据")
    for row_idx in range(50):
        row = TableRow()
        for col_idx in range(5):
            cell = TableCell()
            cell.addElement(P(text=f"cell_{row_idx}_{col_idx}_password"))
            row.addElement(cell)
        table.addElement(row)
    doc.spreadsheet.addElement(table)  # type: ignore[attr-defined]
    buf = io.BytesIO()
    doc.write(buf)
    return buf.getvalue()


def _make_wps_sample() -> bytes:
    """生成典型 WPS 文档样本（OOXML 兼容，本质是 DOCX）。"""
    return _make_docx_sample()


def _make_xls_sample() -> bytes:
    """生成典型 XLS 工作簿样本。

    使用 xlrd 读取测试，但 xlrd 不支持写入；用 mock 跳过。
    """
    pytest.skip("xlrd 不支持写入 XLS，跳过 XLS 基准测试")


def _make_doc_sample() -> bytes:
    """生成典型 DOC 文档样本。

    DOC 为二进制 OLE 格式，难以程序化生成；用 mock 跳过。
    """
    pytest.skip("DOC 为二进制 OLE 格式，难以程序化生成，跳过 DOC 基准测试")


def _make_ppt_sample() -> bytes:
    """生成典型 PPT 演示文稿样本。

    PPT 为二进制 OLE 格式，难以程序化生成；用 mock 跳过。
    """
    pytest.skip("PPT 为二进制 OLE 格式，难以程序化生成，跳过 PPT 基准测试")


def _make_msg_sample() -> bytes:
    """生成典型 MSG 邮件样本。

    MSG 为 Outlook 私有格式，难以程序化生成；用 mock 跳过。
    """
    pytest.skip("MSG 为 Outlook 私有格式，难以程序化生成，跳过 MSG 基准测试")


# ----------------------------- T1 极速：纯文本解码 -----------------------------


@pytest.mark.slow
class TestTier1VeryFast:
    """T1 极速档次基准测试：纯字节解码，无第三方库。

    覆盖纯文本/源代码/配置文件/标记与数据/样式表 5 个子提取器。
    """

    def test_plain_text_extractor_tier(self) -> None:
        """PlainTextExtractor 声明为 T1 极速。"""
        extractor = PlainTextExtractor()
        _assert_tier(extractor, SpeedTier.VERY_FAST)

    def test_source_code_extractor_tier(self) -> None:
        """SourceCodeExtractor 声明为 T1 极速。"""
        extractor = SourceCodeExtractor()
        _assert_tier(extractor, SpeedTier.VERY_FAST)

    def test_text_extraction_speed(self) -> None:
        """100KB 文本提取应在 500ms 内完成（T1 极速基准）。"""
        extractor = PlainTextExtractor()
        data = _make_text_sample(100)
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.VERY_FAST, "PlainTextExtractor")
        assert "password" in extractor.extract_from_bytes(data)


# ----------------------------- T2 快速：标准库解析 -----------------------------


@pytest.mark.slow
class TestTier2Fast:
    """T2 快速档次基准测试：标准库解析。"""

    def test_eml_extractor_tier(self) -> None:
        """EmlExtractor 声明为 T2 快速。"""
        extractor = EmlExtractor()
        _assert_tier(extractor, SpeedTier.FAST)

    def test_eml_extraction_speed(self) -> None:
        """典型 EML 邮件提取应在 1s 内完成（T2 快速基准）。"""
        extractor = EmlExtractor()
        data = _make_eml_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.FAST, "EmlExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content


# ----------------------------- T3 中速：XML 解析 + 树遍历 -----------------------------


@pytest.mark.slow
class TestTier3Medium:
    """T3 中速档次基准测试：单次 XML 解析 + 树遍历。"""

    def test_docx_extractor_tier(self) -> None:
        """DocxExtractor 声明为 T3 中速。"""
        extractor = DocxExtractor()
        _assert_tier(extractor, SpeedTier.MEDIUM)

    def test_docx_extraction_speed(self) -> None:
        """典型 DOCX 文档提取应在 2s 内完成（T3 中速基准）。"""
        extractor = DocxExtractor()
        data = _make_docx_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.MEDIUM, "DocxExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_odt_extractor_tier(self) -> None:
        """OdtExtractor 声明为 T3 中速。"""
        extractor = OdtExtractor()
        _assert_tier(extractor, SpeedTier.MEDIUM)

    def test_odt_extraction_speed(self) -> None:
        """典型 ODT 文档提取应在 2s 内完成（T3 中速基准）。"""
        extractor = OdtExtractor()
        data = _make_odt_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.MEDIUM, "OdtExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_rtf_extractor_tier(self) -> None:
        """RtfExtractor 声明为 T3 中速。"""
        extractor = RtfExtractor()
        _assert_tier(extractor, SpeedTier.MEDIUM)

    def test_rtf_extraction_speed(self) -> None:
        """典型 RTF 文档提取应在 2s 内完成（T3 中速基准）。"""
        extractor = RtfExtractor()
        data = _make_rtf_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.MEDIUM, "RtfExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_wps_extractor_tier(self) -> None:
        """WpsExtractor 声明为 T3 中速。"""
        extractor = WpsExtractor()
        _assert_tier(extractor, SpeedTier.MEDIUM)

    def test_wps_extraction_speed(self) -> None:
        """典型 WPS 文档（OOXML 兼容）提取应在 2s 内完成（T3 中速基准）。"""
        extractor = WpsExtractor()
        data = _make_wps_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.MEDIUM, "WpsExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_msg_extractor_tier(self) -> None:
        """MsgExtractor 声明为 T3 中速。"""
        extractor = MsgExtractor()
        _assert_tier(extractor, SpeedTier.MEDIUM)
        # MSG 样本难以程序化生成，仅验证档次声明


# ----------------------------- T4 慢速：单元格遍历/字节扫描 -----------------------------


@pytest.mark.slow
class TestTier4Slow:
    """T4 慢速档次基准测试：单元格遍历或字节级扫描。"""

    def test_xlsx_extractor_tier(self) -> None:
        """XlsxExtractor 声明为 T4 慢速。"""
        extractor = XlsxExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)

    def test_xlsx_extraction_speed(self) -> None:
        """典型 XLSX 工作簿（100 行 × 10 列）提取应在 5s 内完成（T4 慢速基准）。"""
        extractor = XlsxExtractor()
        data = _make_xlsx_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.SLOW, "XlsxExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_pptx_extractor_tier(self) -> None:
        """PptxExtractor 声明为 T4 慢速。"""
        extractor = PptxExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)

    def test_pptx_extraction_speed(self) -> None:
        """典型 PPTX 演示文稿（5 张幻灯片）提取应在 5s 内完成（T4 慢速基准）。"""
        extractor = PptxExtractor()
        data = _make_pptx_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.SLOW, "PptxExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_ods_extractor_tier(self) -> None:
        """OdsExtractor 声明为 T4 慢速。"""
        extractor = OdsExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)

    def test_ods_extraction_speed(self) -> None:
        """典型 ODS 表格提取应在 5s 内完成（T4 慢速基准）。"""
        extractor = OdsExtractor()
        data = _make_ods_sample()
        elapsed = _measure(extractor.extract_from_bytes, data)
        _assert_time_within_tier(elapsed, SpeedTier.SLOW, "OdsExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content

    def test_xls_extractor_tier(self) -> None:
        """XlsExtractor 声明为 T4 慢速（仅档次声明，样本无法程序化生成）。"""
        extractor = XlsExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)

    def test_doc_extractor_tier(self) -> None:
        """DocExtractor 声明为 T4 慢速（仅档次声明，样本无法程序化生成）。"""
        from fuscan.extractors.legacy_office import DocExtractor

        extractor = DocExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)

    def test_ppt_extractor_tier(self) -> None:
        """PptExtractor 声明为 T4 慢速（仅档次声明，样本无法程序化生成）。"""
        from fuscan.extractors.legacy_office import PptExtractor

        extractor = PptExtractor()
        _assert_tier(extractor, SpeedTier.SLOW)


# ----------------------------- T5 极慢：复杂布局分析 -----------------------------


@pytest.mark.slow
class TestTier5VerySlow:
    """T5 极慢档次基准测试：复杂页面布局分析。"""

    def test_pdf_extractor_tier(self) -> None:
        """PdfExtractor 声明为 T5 极慢。"""
        extractor = PdfExtractor()
        _assert_tier(extractor, SpeedTier.VERY_SLOW)

    def test_pdf_extraction_speed(self) -> None:
        """典型 PDF 文档（单页 20 行）提取应在 10s 内完成（T5 极慢基准）。"""
        extractor = PdfExtractor()
        data = _make_pdf_sample()
        elapsed = _measure(extractor.extract_from_bytes, data, iterations=1)
        _assert_time_within_tier(elapsed, SpeedTier.VERY_SLOW, "PdfExtractor")
        content = extractor.extract_from_bytes(data)
        assert "password" in content


# ----------------------------- 档次完整性校验 -----------------------------


class TestSpeedTierCompleteness:
    """验证所有注册提取器都声明了有效的 speed_tier（非 slow 测试）。"""

    def test_all_extractors_have_valid_speed_tier(self) -> None:
        """default_registry 中所有提取器都应声明有效的 speed_tier。"""
        from fuscan.extractors import default_registry

        for _ext, extractor in default_registry._extractors.items():  # type: ignore[attr-defined]
            tier = extractor.speed_tier
            assert isinstance(tier, SpeedTier), f"{type(extractor).__name__} 返回非 SpeedTier 值"
            assert tier in SpeedTier, f"{type(extractor).__name__} 返回无效档次 {tier}"

    def test_speed_tier_labels_are_chinese(self) -> None:
        """SpeedTier.label 返回中文标签（T1-T5 + 中文描述）。"""
        assert SpeedTier.VERY_FAST.label == "T1 极速"
        assert SpeedTier.FAST.label == "T2 快速"
        assert SpeedTier.MEDIUM.label == "T3 中速"
        assert SpeedTier.SLOW.label == "T4 慢速"
        assert SpeedTier.VERY_SLOW.label == "T5 极慢"

    def test_speed_tier_descriptions_are_chinese(self) -> None:
        """SpeedTier.description 返回中文说明（含耗时范围）。"""
        assert "< 10ms/MB" in SpeedTier.VERY_FAST.description
        assert "10-50ms/MB" in SpeedTier.FAST.description
        assert "50-200ms/MB" in SpeedTier.MEDIUM.description
        assert "200-1000ms/MB" in SpeedTier.SLOW.description
        assert "> 1000ms/MB" in SpeedTier.VERY_SLOW.description

    def test_speed_tier_colors_from_green_to_red(self) -> None:
        """SpeedTier.color 返回从绿到红的十六进制色值（iter-91）。"""
        assert SpeedTier.VERY_FAST.color == "#28A745"  # 绿
        assert SpeedTier.FAST.color == "#17A2B8"  # 青
        assert SpeedTier.MEDIUM.color == "#FFC107"  # 琥珀
        assert SpeedTier.SLOW.color == "#FD7E14"  # 橙
        assert SpeedTier.VERY_SLOW.color == "#DC3545"  # 红
