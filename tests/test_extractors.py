"""提取器单元测试。

使用对应库动态生成测试 fixture 文件，避免二进制 fixture 入仓。
PDF/ODT/ODS 等较难动态生成的格式，使用 mock 或跳过。
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest

from fuscan.extractors import (
    DocExtractor,
    DocxExtractor,
    EmlExtractor,
    ExtractorError,
    ExtractorRegistry,
    MsgExtractor,
    OdtExtractor,
    PdfExtractor,
    PptExtractor,
    PptxExtractor,
    RtfExtractor,
    TextExtractor,
    WpsExtractor,
    XlsExtractor,
    XlsxExtractor,
    clear_content_cache,
    default_registry,
    extract_content,
    extract_content_cached,
    extract_content_from_bytes,
    extract_content_with_fallback,
    get_extractor,
)
from fuscan.extractors.base import SpeedTier
from fuscan.extractors.spreadsheet import OdsExtractor


def _make_ooxml_zip(entry_name: str, content: str = "fake") -> bytes:
    """创建包含指定条目的有效 ZIP，用于测试 OOXML 类型检测。

    :param entry_name: ZIP 内部条目名（如 ``word/document.xml``）
    :param content: 条目内容（默认 ``"fake"``，用于触发解析失败）
    :return: ZIP 文件字节
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(entry_name, content)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Fixture 工厂
# ---------------------------------------------------------------------------


@pytest.fixture()
def text_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.txt"
    path.write_text("hello password world\n第二行内容\n", encoding="utf-8")
    return path


@pytest.fixture()
def gbk_file(tmp_path: Path) -> Path:
    path = tmp_path / "gbk.txt"
    path.write_bytes("这是一个包含密码字段的配置文件，密码为 password123，请妥善保管。PASSWORD".encode("gbk"))
    return path


@pytest.fixture()
def empty_file(tmp_path: Path) -> Path:
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")
    return path


@pytest.fixture()
def docx_file(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_paragraph("段落一 含 password")
    doc.add_paragraph("段落二 正常内容")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "密码"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "pwd123"
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


@pytest.fixture()
def pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "标题 含 secret"  # pyrefly: ignore [missing-attribute]
    slide.placeholders[1].text = "幻灯片内容"  # pyrefly: ignore [missing-attribute]
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def xlsx_file(tmp_path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "数据"
    ws["A1"] = "姓名"
    ws["B1"] = "密码"
    ws["A2"] = "张三"
    ws["B2"] = "pwd123"
    wb.create_sheet("空表")
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# TextExtractor
# ---------------------------------------------------------------------------


class TestTextExtractor:
    def test_extract_utf8(self, text_file: Path) -> None:
        extractor = TextExtractor()
        content = extractor.extract(text_file)
        assert "hello password world" in content
        assert "第二行内容" in content

    def test_extract_empty(self, empty_file: Path) -> None:
        extractor = TextExtractor()
        assert extractor.extract(empty_file) == ""

    def test_supported_extensions(self) -> None:
        extractor = TextExtractor()
        assert "txt" in extractor.supported_extensions
        assert "md" in extractor.supported_extensions
        assert "py" in extractor.supported_extensions

    def test_max_size_limit(self, tmp_path: Path) -> None:
        path = tmp_path / "big.txt"
        path.write_text("x" * 100, encoding="utf-8")
        extractor = TextExtractor(max_size=10)
        assert extractor.extract(path) == ""

    def test_nonexistent_file_raises(self, tmp_path: Path) -> None:
        extractor = TextExtractor()
        with pytest.raises(ExtractorError, match="无法读取文件大小"):
            extractor.extract(tmp_path / "missing.txt")

    def test_charset_normalizer_fallback(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 未安装时回退到 UTF-8/GBK 解码。"""
        path = tmp_path / "fallback.txt"
        path.write_text("回退解码 password", encoding="utf-8")

        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "charset_normalizer":
                raise ImportError("No module named 'charset_normalizer'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)

        content = TextExtractor().extract(path)
        assert "回退解码 password" in content

    def test_gbk_decoding(self, gbk_file: Path) -> None:
        """GBK 编码文件应能正确解码。"""
        extractor = TextExtractor()
        content = extractor.extract(gbk_file)
        assert "密码" in content
        assert "password123" in content

    def test_read_bytes_os_error(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """read_bytes 失败时抛出 ExtractorError。"""
        path = tmp_path / "test.txt"
        path.write_text("hello", encoding="utf-8")

        original_read_bytes = Path.read_bytes

        def mock_read_bytes(self: Path) -> bytes:
            if self == path:
                raise OSError("模拟读取失败")
            return original_read_bytes(self)

        monkeypatch.setattr(Path, "read_bytes", mock_read_bytes)
        with pytest.raises(ExtractorError, match="文件读取失败"):
            TextExtractor().extract(path)

    def test_charset_normalizer_exception_fallback(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 抛异常时回退到 UTF-8 解码。"""
        path = tmp_path / "test.txt"
        path.write_text("异常回退 password", encoding="utf-8")

        def fake_from_bytes(data: bytes):
            raise RuntimeError("模拟检测异常")

        monkeypatch.setattr("charset_normalizer.from_bytes", fake_from_bytes)
        content = TextExtractor().extract(path)
        assert "异常回退 password" in content

    def test_charset_normalizer_none_fallback_to_latin1(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 返回 None 时回退到 latin-1 解码任意字节。"""
        path = tmp_path / "test.txt"
        path.write_bytes(b"\xff\xfe\xfd")

        monkeypatch.setattr("charset_normalizer.from_bytes", lambda data: type("R", (), {"best": lambda self: None})())
        content = TextExtractor().extract(path)
        assert isinstance(content, str)
        # latin-1 能解码任意字节，不应为空
        assert len(content) == 3

    def test_normalizes_crlf_to_lf(self, tmp_path: Path) -> None:
        """CRLF 行尾应规范化为 LF，保证跨平台 CONTENT EQUALS 比较一致。"""
        path = tmp_path / "crlf.txt"
        path.write_bytes(b"line1\r\nline2\r\n")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"
        assert "\r\n" not in content

    def test_normalizes_cr_to_lf(self, tmp_path: Path) -> None:
        """旧式 CR 行尾应规范化为 LF。"""
        path = tmp_path / "cr.txt"
        path.write_bytes(b"line1\rline2\r")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"

    def test_lf_preserved(self, tmp_path: Path) -> None:
        """LF 行尾保持不变。"""
        path = tmp_path / "lf.txt"
        path.write_bytes(b"line1\nline2\n")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"


# ---------------------------------------------------------------------------
# DocxExtractor
# ---------------------------------------------------------------------------


class TestDocxExtractor:
    def test_extract_paragraphs_and_table(self, docx_file: Path) -> None:
        extractor = DocxExtractor()
        content = extractor.extract(docx_file)
        assert "段落一 含 password" in content
        assert "段落二 正常内容" in content
        assert "姓名" in content
        assert "pwd123" in content

    def test_supported_extensions(self) -> None:
        assert DocxExtractor().supported_extensions == ("docx",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.docx"
        path.write_text("not a docx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="DOCX 解析失败"):
            DocxExtractor().extract(path)


# ---------------------------------------------------------------------------
# PptxExtractor
# ---------------------------------------------------------------------------


class TestPptxExtractor:
    def test_extract_slide_text(self, pptx_file: Path) -> None:
        extractor = PptxExtractor()
        content = extractor.extract(pptx_file)
        assert "标题 含 secret" in content
        assert "幻灯片内容" in content

    def test_supported_extensions(self) -> None:
        assert PptxExtractor().supported_extensions == ("pptx",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.pptx"
        path.write_text("not a pptx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="PPTX 解析失败"):
            PptxExtractor().extract(path)

    def test_pptx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-pptx 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.pptx"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-pptx 未安装"):
            PptxExtractor().extract(path)

    def test_pptx_with_table_and_notes(self, tmp_path: Path) -> None:
        """PPTX 含表格和备注时应提取这些内容。"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "表格测试"  # pyrefly: ignore [missing-attribute]
        # 添加表格
        table_shape = slide.shapes.add_table(
            rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(1)
        )
        table = table_shape.table
        table.cell(0, 0).text = "键"
        table.cell(0, 1).text = "密码"
        table.cell(1, 0).text = "user"
        table.cell(1, 1).text = "pwd123"
        # 添加备注
        slide.notes_slide.notes_text_frame.text = "备注内容 secret"  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "table_notes.pptx"
        prs.save(str(path))

        content = PptxExtractor().extract(path)
        assert "表格测试" in content
        assert "pwd123" in content
        assert "[备注]" in content
        assert "备注内容 secret" in content

    def test_pptx_empty_slide_skipped(self, tmp_path: Path) -> None:
        """空幻灯片应被跳过（不添加 --- 幻灯片 --- 分隔符）。"""
        from pptx import Presentation

        prs = Presentation()
        # 添加一个空布局的幻灯片（无文本）
        prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        path = tmp_path / "empty.pptx"
        prs.save(str(path))

        content = PptxExtractor().extract(path)
        assert "幻灯片" not in content


class TestDocxExtractorExtra:
    """DocxExtractor 额外覆盖。"""

    def test_docx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-docx 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.docx"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "docx":
                raise ImportError("No module named 'docx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-docx 未安装"):
            DocxExtractor().extract(path)

    def test_docx_with_header_footer(self, tmp_path: Path) -> None:
        """DOCX 含页眉页脚时应提取这些内容。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("正文 password")
        # 添加页眉页脚
        section = doc.sections[0]
        section.header.paragraphs[0].text = "页眉内容 secret"
        section.footer.paragraphs[0].text = "页脚信息"
        path = tmp_path / "header_footer.docx"
        doc.save(str(path))

        content = DocxExtractor().extract(path)
        assert "正文 password" in content
        assert "页眉内容 secret" in content
        assert "页脚信息" in content


# ---------------------------------------------------------------------------
# XlsxExtractor
# ---------------------------------------------------------------------------


class TestXlsxExtractor:
    def test_extract_cells(self, xlsx_file: Path) -> None:
        extractor = XlsxExtractor()
        content = extractor.extract(xlsx_file)
        assert "姓名" in content
        assert "pwd123" in content
        assert "数据" in content  # 工作表名

    def test_supported_extensions(self) -> None:
        exts = XlsxExtractor().supported_extensions
        assert "xlsx" in exts
        assert "xlsm" in exts

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.xlsx"
        path.write_text("not xlsx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="XLSX 解析失败"):
            XlsxExtractor().extract(path)

    def test_max_rows_limit(self, xlsx_file: Path) -> None:
        extractor = XlsxExtractor(max_rows=1)
        content = extractor.extract(xlsx_file)
        # 只读了 1 行，应该有表头但无数据行
        assert "姓名" in content

    def test_max_cols_limit(self, xlsx_file: Path) -> None:
        """列数超过上限时截断。"""
        extractor = XlsxExtractor(max_cols=1)
        content = extractor.extract(xlsx_file)
        # 只有第 1 列，应包含姓名但不包含密码列
        assert "姓名" in content

    def test_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.xlsx"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            XlsxExtractor().extract(path)


# ---------------------------------------------------------------------------
# WpsExtractor
# ---------------------------------------------------------------------------


class TestWpsExtractor:
    def test_supported_extensions(self) -> None:
        exts = WpsExtractor().supported_extensions
        assert "wps" in exts
        assert "et" in exts
        assert "dps" in exts

    def test_non_ooxml_returns_empty(self, tmp_path: Path) -> None:
        """旧版二进制 WPS 格式应返回空字符串。"""
        path = tmp_path / "old.wps"
        path.write_bytes(b"\xd0\xcf\x11\xe0 old binary format")
        assert WpsExtractor().extract(path) == ""

    def test_ooxml_wps_text(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .wps 文件应能提取文本（实际是 DOCX 内容）。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("wps 内容 password")
        path = tmp_path / "test.wps"
        doc.save(str(path))

        content = WpsExtractor().extract(path)
        assert "wps 内容 password" in content

    def test_ooxml_et_sheet(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .et 文件应能提取表格内容。"""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws["A1"] = "et_password"
        path = tmp_path / "test.et"
        wb.save(str(path))

        content = WpsExtractor().extract(path)
        assert "et_password" in content

    def test_ooxml_dps_slides(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .dps 文件应能提取演示内容。"""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "dps 标题 password"  # pyrefly: ignore [missing-attribute]
        slide.placeholders[1].text = "幻灯片内容"  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "test.dps"
        prs.save(str(path))

        content = WpsExtractor().extract(path)
        assert "dps 标题 password" in content
        assert "幻灯片内容" in content

    def test_wps_docx_with_table(self, tmp_path: Path) -> None:
        """WPS 文字文档含表格时应提取表格内容。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("正文 password")
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "键"
        table.cell(0, 1).text = "值 secret"
        path = tmp_path / "table.wps"
        doc.save(str(path))

        content = WpsExtractor().extract(path)
        assert "正文 password" in content
        assert "键" in content
        assert "值 secret" in content

    def test_wps_invalid_docx_raises(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 docx 内容损坏时应抛出 ExtractorError。"""
        path = tmp_path / "bad.wps"
        path.write_bytes(_make_ooxml_zip("word/document.xml", "corrupt xml"))
        with pytest.raises(ExtractorError, match="WPS 文字文档解析失败"):
            WpsExtractor().extract(path)

    def test_wps_invalid_et_raises(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 xlsx 内容损坏时应抛出 ExtractorError。"""
        path = tmp_path / "bad.et"
        path.write_bytes(_make_ooxml_zip("xl/workbook.xml", "corrupt xml"))
        with pytest.raises(ExtractorError, match="WPS 表格 解析失败"):
            WpsExtractor().extract(path)

    def test_wps_invalid_dps_raises(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 pptx 内容损坏时应抛出 ExtractorError。"""
        path = tmp_path / "bad.dps"
        path.write_bytes(_make_ooxml_zip("ppt/presentation.xml", "corrupt xml"))
        with pytest.raises(ExtractorError, match="WPS 演示解析失败"):
            WpsExtractor().extract(path)

    def test_wps_corrupt_zip_returns_empty(self, tmp_path: Path) -> None:
        """ZIP 头存在但 ZIP 损坏时应返回空字符串（无法判定子类型）。"""
        path = tmp_path / "corrupt.wps"
        path.write_bytes(b"PK\x03\x04 corrupted content")
        assert WpsExtractor().extract(path) == ""

    def test_wps_detect_ooxml_type_bad_zip_returns_none(self) -> None:
        """_detect_ooxml_type 对损坏的 ZIP 数据返回 None。"""
        assert WpsExtractor()._detect_ooxml_type(b"PK\x03\x04 corrupted") is None

    def test_wps_detect_ooxml_type_unknown_returns_none(self) -> None:
        """_detect_ooxml_type 对未知 OOXML 子类型返回 None。"""
        data = _make_ooxml_zip("unknown/entry.xml")
        assert WpsExtractor()._detect_ooxml_type(data) is None

    def test_wps_detect_ooxml_type_docx(self) -> None:
        """_detect_ooxml_type 正确识别 docx 子类型。"""
        data = _make_ooxml_zip("word/document.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "docx"

    def test_wps_detect_ooxml_type_xlsx(self) -> None:
        """_detect_ooxml_type 正确识别 xlsx 子类型。"""
        data = _make_ooxml_zip("xl/workbook.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "xlsx"

    def test_wps_detect_ooxml_type_pptx(self) -> None:
        """_detect_ooxml_type 正确识别 pptx 子类型。"""
        data = _make_ooxml_zip("ppt/presentation.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "pptx"


class TestWpsExtractorErrorPaths:
    """WPS 提取器异常路径覆盖。"""

    def test_extract_corrupt_zip_returns_empty(self, tmp_path: Path) -> None:
        """ZIP 头存在但内容损坏时应返回空字符串（无法判定子类型）。"""
        path = tmp_path / "file.unknown"
        path.write_bytes(b"PK\x03\x04 fake content")
        # ZIP 头存在但非有效 ZIP，_detect_ooxml_type 返回 None，走 return ""
        extractor = WpsExtractor()
        result = extractor.extract(path)
        assert result == ""

    def test_extract_as_docx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-docx 未安装时 _extract_as_docx 应抛出 ExtractorError。"""
        path = tmp_path / "test.wps"
        # 用有效 ZIP（含 word/document.xml）让类型检测返回 docx，触发后续 import
        path.write_bytes(_make_ooxml_zip("word/document.xml"))
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "docx":
                raise ImportError("No module named 'docx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-docx 未安装"):
            WpsExtractor().extract(path)

    def test_extract_as_xlsx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时 _extract_as_xlsx 应抛出 ExtractorError。"""
        path = tmp_path / "test.et"
        path.write_bytes(_make_ooxml_zip("xl/workbook.xml"))
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            WpsExtractor().extract(path)

    def test_extract_as_pptx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-pptx 未安装时 _extract_as_pptx 应抛出 ExtractorError。"""
        path = tmp_path / "test.dps"
        path.write_bytes(_make_ooxml_zip("ppt/presentation.xml"))
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-pptx 未安装"):
            WpsExtractor().extract(path)


# ---------------------------------------------------------------------------
# PdfExtractor（依赖 pypdf，使用 mock 避免真实 PDF）
# iter-91：pdf_oxide 可用时优先走 Rust 后端，以下测试强制走 pypdf 回退路径
# ---------------------------------------------------------------------------


class TestPdfExtractor:
    """PDF 提取器测试。

    iter-91 起 PdfExtractor 优先使用 pdf_oxide（Rust + PyO3），回退到 pypdf。
    以下 mock 测试通过 ``monkeypatch`` 强制 ``_PDF_OXIDE_AVAILABLE = False``，
    验证 pypdf 回退路径的正确性。pdf_oxide 后端的正确性由基准测试覆盖。
    """

    @pytest.fixture(autouse=True)
    def _force_pypdf_fallback(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """强制走 pypdf 回退路径，绕过 pdf_oxide（iter-91）。"""
        monkeypatch.setattr("fuscan.extractors.pdf._PDF_OXIDE_AVAILABLE", False)

    def test_supported_extensions(self) -> None:
        assert PdfExtractor().supported_extensions == ("pdf",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.pdf"
        path.write_text("not a pdf", encoding="utf-8")
        with pytest.raises(ExtractorError):
            PdfExtractor().extract(path)

    def test_extract_with_mock_reader(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """使用 mock PdfReader 测试页面提取。"""
        path = tmp_path / "fake.pdf"
        path.write_bytes(b"fake pdf content")

        class FakePage:
            def extract_text(self) -> str:
                return "页面文本含 password"

        class FakeReader:
            def __init__(self) -> None:
                self.is_encrypted = False
                self.pages = [FakePage(), FakePage()]

        class FakePdfModule:
            PdfReader = staticmethod(lambda _: FakeReader())

            class errors:
                class PdfReadError(Exception):
                    pass

        import sys

        monkeypatch.setitem(sys.modules, "pypdf", FakePdfModule)
        monkeypatch.setitem(sys.modules, "pypdf.errors", type("errors", (), {"PdfReadError": Exception}))

        content = PdfExtractor().extract(path)
        assert "页面文本含 password" in content

    def test_encrypted_pdf_returns_empty(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """加密 PDF 应返回空字符串。"""
        path = tmp_path / "encrypted.pdf"
        path.write_bytes(b"encrypted")

        class FakeReader:
            def __init__(self) -> None:
                self.is_encrypted = True
                self.pages = []

        import sys

        fake_module = type("pypdf", (), {"PdfReader": staticmethod(lambda _: FakeReader())})
        fake_errors = type("errors", (), {"PdfReadError": Exception})
        monkeypatch.setitem(sys.modules, "pypdf", fake_module)
        monkeypatch.setitem(sys.modules, "pypdf.errors", fake_errors)

        assert PdfExtractor().extract(path) == ""

    def test_page_extraction_error_continues(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """单页提取失败应跳过该页继续处理。"""
        path = tmp_path / "mixed.pdf"
        path.write_bytes(b"mixed")

        class GoodPage:
            def extract_text(self) -> str:
                return "正常页面"

        class BadPage:
            def extract_text(self) -> str:
                raise RuntimeError("解析失败")

        class FakeReader:
            def __init__(self) -> None:
                self.is_encrypted = False
                self.pages = [BadPage(), GoodPage()]

        fake_module = type("pypdf", (), {"PdfReader": staticmethod(lambda _: FakeReader())})
        fake_errors = type("errors", (), {"PdfReadError": Exception})
        import sys

        monkeypatch.setitem(sys.modules, "pypdf", fake_module)
        monkeypatch.setitem(sys.modules, "pypdf.errors", fake_errors)

        content = PdfExtractor().extract(path)
        assert "正常页面" in content

    def test_empty_page_text_skipped(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract_text 返回空字符串的页面应被跳过（不加入 parts）。"""
        path = tmp_path / "empty_pages.pdf"
        path.write_bytes(b"empty")

        class EmptyPage:
            def extract_text(self) -> str:
                return ""

        class GoodPage:
            def extract_text(self) -> str:
                return "有内容 password"

        class FakeReader:
            def __init__(self) -> None:
                self.is_encrypted = False
                self.pages = [EmptyPage(), GoodPage(), EmptyPage()]

        fake_module = type("pypdf", (), {"PdfReader": staticmethod(lambda _: FakeReader())})
        fake_errors = type("errors", (), {"PdfReadError": Exception})
        import sys

        monkeypatch.setitem(sys.modules, "pypdf", fake_module)
        monkeypatch.setitem(sys.modules, "pypdf.errors", fake_errors)

        content = PdfExtractor().extract(path)
        assert "有内容 password" in content
        assert content == "有内容 password"

    def test_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """pypdf 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.pdf"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "pypdf":
                raise ImportError("No module named 'pypdf'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="pypdf 未安装"):
            PdfExtractor().extract(path)

    def test_pdf_open_generic_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """PdfReader 抛出非 PdfReadError 异常时应包装为 ExtractorError。"""
        path = tmp_path / "corrupt.pdf"
        path.write_bytes(b"corrupt")

        def raise_oserror(_path: str) -> object:
            raise OSError("模拟打开失败")

        class FakeErrors:
            class PdfReadError(Exception):
                pass

        fake_module = type("pypdf", (), {"PdfReader": staticmethod(raise_oserror)})
        import sys

        monkeypatch.setitem(sys.modules, "pypdf", fake_module)
        monkeypatch.setitem(sys.modules, "pypdf.errors", FakeErrors)

        with pytest.raises(ExtractorError, match="PDF 打开失败"):
            PdfExtractor().extract(path)

    def test_pdf_read_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """PdfReader 抛出 PdfReadError 时应包装为 ExtractorError。"""
        path = tmp_path / "bad.pdf"
        path.write_bytes(b"bad")

        class FakePdfReadError(Exception):
            pass

        def raise_pdf_error(_path: str) -> object:
            raise FakePdfReadError("模拟解析失败")

        fake_module = type("pypdf", (), {"PdfReader": staticmethod(raise_pdf_error)})
        import sys

        monkeypatch.setitem(sys.modules, "pypdf", fake_module)
        monkeypatch.setitem(sys.modules, "pypdf.errors", type("errors", (), {"PdfReadError": FakePdfReadError}))

        with pytest.raises(ExtractorError, match="PDF 解析失败"):
            PdfExtractor().extract(path)


# ---------------------------------------------------------------------------
# PdfExtractor pdf_oxide 后端测试（iter-91）
# ---------------------------------------------------------------------------


class TestPdfExtractorOxideBackend:
    """pdf_oxide（Rust + PyO3）后端测试。

    仅在 pdf_oxide 已安装时运行，验证真实 PDF 提取与 speed_tier 动态返回。
    不使用 mock，生成真实 PDF 样本验证端到端正确性。
    """

    @pytest.fixture()
    def pdf_sample(self, tmp_path: Path) -> bytes:
        """用 reportlab 生成含 password 关键词的 PDF 样本。"""
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate

        path = tmp_path / "sample.pdf"
        styles = getSampleStyleSheet()
        doc = SimpleDocTemplate(str(path), pagesize=letter)
        doc.build([Paragraph("This document contains a secret password.", styles["Normal"])])
        return path.read_bytes()

    def test_oxide_speed_tier_is_fast_when_available(self) -> None:
        """pdf_oxide 可用时 speed_tier 返回 T2 快速。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        assert PdfExtractor().speed_tier == SpeedTier.FAST

    def test_oxide_extract_real_pdf(self, pdf_sample: bytes) -> None:
        """pdf_oxide 后端提取真实 PDF 应包含 password 关键词。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(pdf_sample)
        assert "password" in content.lower()

    def test_oxide_extract_path_matches_bytes(self, pdf_sample: bytes, tmp_path: Path) -> None:
        """pdf_oxide 从 path 与从 bytes 提取结果一致。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        path = tmp_path / "sample.pdf"
        path.write_bytes(pdf_sample)
        extractor = PdfExtractor()
        assert extractor.extract(path) == extractor.extract_from_bytes(pdf_sample)

    def test_oxide_invalid_bytes_raises(self) -> None:
        """pdf_oxide 后端对无效字节应抛出 ExtractorError。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        with pytest.raises(ExtractorError, match="PDF 解析失败"):
            PdfExtractor().extract_from_bytes(b"not a pdf")

    def test_oxide_empty_bytes_raises(self) -> None:
        """pdf_oxide 后端对空字节应抛出 ExtractorError（非加密异常）。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        with pytest.raises(ExtractorError):
            PdfExtractor().extract_from_bytes(b"")

    def test_oxide_extract_missing_file_raises(self, tmp_path: Path) -> None:
        """pdf_oxide 后端 extract() 读取缺失文件应抛出 ExtractorError。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        with pytest.raises(ExtractorError, match="文件读取失败"):
            PdfExtractor().extract(tmp_path / "missing.pdf")

    def test_oxide_returns_empty_for_empty_pdf(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """pdf_oxide 后端 to_plain_text_all 返回空时返回空字符串。"""
        from fuscan.extractors.pdf import _PDF_OXIDE_AVAILABLE

        if not _PDF_OXIDE_AVAILABLE:
            pytest.skip("pdf_oxide 未安装")
        extractor = PdfExtractor()

        class FakeDoc:
            @staticmethod
            def to_plain_text_all() -> str:
                return ""

        monkeypatch.setattr(
            "fuscan.extractors.pdf._PdfOxideDocument.from_bytes",
            lambda data: FakeDoc(),
        )
        assert extractor.extract_from_bytes(b"fake but callable") == ""


# ---------------------------------------------------------------------------
# OdtExtractor / OdsExtractor（依赖 odfpy）
# ---------------------------------------------------------------------------


class TestOdfExtractors:
    def test_odt_supported_extensions(self) -> None:
        assert OdtExtractor().supported_extensions == ("odt",)

    def test_ods_supported_extensions(self) -> None:
        assert OdsExtractor().supported_extensions == ("ods",)

    def test_odt_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.odt"
        path.write_text("not odt", encoding="utf-8")
        with pytest.raises(ExtractorError, match="ODT 解析失败"):
            OdtExtractor().extract(path)

    def test_odt_extract_real_file(self, tmp_path: Path) -> None:
        """使用 odfpy 创建真实 ODT 文件并提取。"""
        from odf.opendocument import OpenDocumentText
        from odf.text import H, P

        doc = OpenDocumentText()
        p = P(text="段落含 password 内容")
        doc.text.addElement(p)  # pyrefly: ignore [missing-attribute]
        h = H(outlinelevel="1", text="标题 secret")
        doc.text.addElement(h)  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "real.odt"
        doc.save(str(path))

        content = OdtExtractor().extract(path)
        assert "password" in content
        assert "secret" in content

    def test_ods_extract_real_file(self, tmp_path: Path) -> None:
        """使用 odfpy 创建真实 ODS 文件并提取。"""
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableCell, TableRow
        from odf.text import P

        doc = OpenDocumentSpreadsheet()
        table = Table(name="数据")
        row = TableRow()
        cell = TableCell()
        cell.addElement(P(text="cell_password"))
        row.addElement(cell)
        table.addElement(row)
        doc.spreadsheet.addElement(table)  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "real.ods"
        doc.save(str(path))

        content = OdsExtractor().extract(path)
        assert "cell_password" in content

    def test_ods_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.ods"
        path.write_text("not ods", encoding="utf-8")
        with pytest.raises(ExtractorError, match="ODS 解析失败"):
            OdsExtractor().extract(path)

    def test_odt_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """odfpy 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.odt"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "odf.opendocument":
                raise ImportError("No module named 'odf'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="odfpy 未安装"):
            OdtExtractor().extract(path)

    def test_ods_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """odfpy 未安装时 OdsExtractor 应抛出 ExtractorError。"""
        path = tmp_path / "test.ods"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "odf.opendocument":
                raise ImportError("No module named 'odf'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="odfpy 未安装"):
            OdsExtractor().extract(path)

    def test_ods_cell_text_exception_returns_empty(self) -> None:
        """_extract_cell_text 在 str(cell) 抛异常时应返回空字符串。"""
        extractor = OdsExtractor()

        class BadCell:
            def __str__(self) -> str:
                raise RuntimeError("模拟单元格转换失败")

        assert extractor._extract_cell_text(BadCell()) == ""


# ---------------------------------------------------------------------------
# ExtractorRegistry
# ---------------------------------------------------------------------------


class TestExtractorRegistry:
    def test_register_and_get(self) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        assert registry.get("txt") is not None
        assert registry.get("TXT") is not None  # 大小写不敏感
        assert registry.get("missing") is None

    def test_registered_extensions(self) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        exts = registry.registered_extensions
        assert "txt" in exts
        assert "md" in exts

    def test_extract_with_registered(self, text_file: Path) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        content = registry.extract(text_file)
        assert "hello password world" in content

    def test_extract_without_extractor_returns_empty(self, tmp_path: Path) -> None:
        path = tmp_path / "unknown.xyz"
        path.write_text("content", encoding="utf-8")
        registry = ExtractorRegistry()
        assert registry.extract(path) == ""

    def test_default_registry_has_all(self) -> None:
        exts = default_registry.registered_extensions
        for expected in ("txt", "pdf", "docx", "pptx", "xlsx", "odt", "ods", "wps"):
            assert expected in exts, f"默认注册表缺少 {expected}"

    def test_list_extractors_returns_unique_entries(self) -> None:
        """list_extractors 返回去重后的提取器列表（同一实例多扩展名合并为一项）。"""
        extractors = default_registry.list_extractors()
        # 14 个提取器类，每个对应一项
        class_names = [entry[0] for entry in extractors]
        assert len(class_names) == len(set(class_names)), "提取器列表有重复"
        # 按 display_name 排序
        display_names = [entry[1] for entry in extractors]
        assert display_names == sorted(display_names)

    def test_list_extractors_entry_format(self) -> None:
        """list_extractors 返回元组格式为 (class_name, display_name, extensions, speed_tier)。

        iter-90 起新增 speed_tier 字段（SpeedTier 枚举）。
        """
        from fuscan.extractors.base import SpeedTier

        extractors = default_registry.list_extractors()
        for class_name, display_name, exts, tier in extractors:
            assert isinstance(class_name, str) and class_name
            assert isinstance(display_name, str) and display_name
            assert isinstance(exts, tuple) and exts
            assert isinstance(tier, SpeedTier)
            # 扩展名均为小写无点
            for e in exts:
                assert e == e.lower().lstrip(".")

    def test_display_name_returns_chinese(self) -> None:
        """各提取器 display_name 返回非空中文名称。"""
        names = {
            TextExtractor: "纯文本",
            PdfExtractor: "PDF",
            DocxExtractor: "Word（DOCX）",
            PptxExtractor: "PowerPoint（PPTX）",
            XlsxExtractor: "Excel（XLSX）",
            OdsExtractor: "ODS 表格",
            OdtExtractor: "ODT 文档",
            WpsExtractor: "WPS 文档",
            RtfExtractor: "RTF",
            EmlExtractor: "邮件（EML）",
            MsgExtractor: "Outlook 邮件（MSG）",
            XlsExtractor: "Excel（XLS）",
            DocExtractor: "Word（DOC）",
            PptExtractor: "PowerPoint（PPT）",
        }
        for cls, expected in names.items():
            assert cls().display_name == expected, f"{cls.__name__}.display_name 应为 {expected}"


# ---------------------------------------------------------------------------
# 集成函数
# ---------------------------------------------------------------------------


class TestExtractContent:
    def test_extract_text(self, text_file: Path) -> None:
        content = extract_content(text_file)
        assert "hello password world" in content

    def test_extract_docx(self, docx_file: Path) -> None:
        content = extract_content(docx_file)
        assert "段落一 含 password" in content

    def test_extract_unknown_returns_empty(self, tmp_path: Path) -> None:
        path = tmp_path / "unknown.xyz"
        path.write_text("content", encoding="utf-8")
        assert extract_content(path) == ""

    def test_fallback_returns_extracted_content(self, text_file: Path) -> None:
        """提取器成功时返回提取的内容。"""
        content = extract_content_with_fallback(text_file)
        assert "hello password world" in content

    def test_fallback_extractor_failure_falls_back_to_text(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """提取器抛异常时回退到纯文本读取。"""

        def raise_extract(p: Path) -> str:
            raise RuntimeError("提取失败")

        path = tmp_path / "a.txt"
        path.write_text("plain content", encoding="utf-8")
        monkeypatch.setattr("fuscan.extractors.base.extract_content", raise_extract)
        assert extract_content_with_fallback(path) == "plain content"

    def test_fallback_read_text_oserror_propagates(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """纯文本回退读取失败时 OSError 向上传播。"""

        def raise_extract(p: Path) -> str:
            raise RuntimeError("提取失败")

        path = tmp_path / "nonexistent.txt"
        monkeypatch.setattr("fuscan.extractors.base.extract_content", raise_extract)
        with pytest.raises(OSError):
            extract_content_with_fallback(path)

    def test_get_extractor_returns_none_for_unknown(self) -> None:
        assert get_extractor("xyz") is None

    def test_get_extractor_returns_instance(self) -> None:
        extractor = get_extractor("txt")
        assert extractor is not None
        assert isinstance(extractor, TextExtractor)


# ---------------------------------------------------------------------------
# Scanner 集成
# ---------------------------------------------------------------------------


class TestScannerWithExtractors:
    def test_scan_docx_content(self, docx_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.CRITICAL,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="password"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(docx_file)
        assert result.has_hit
        assert result.hits[0].rule_name == "敏感词"

    def test_scan_xlsx_content(self, xlsx_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.WARNING,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="pwd123"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(xlsx_file)
        assert result.has_hit


# ---------------------------------------------------------------------------
# extract_from_bytes：各提取器从内存字节提取（消除双重 I/O）
# ---------------------------------------------------------------------------


class TestExtractFromBytes:
    """各提取器 extract_from_bytes 与 extract(path) 结果一致性。"""

    def test_text_extract_from_bytes_matches_path(self, text_file: Path) -> None:
        """TextExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = text_file.read_bytes()
        extractor = TextExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(text_file)

    def test_text_extract_from_bytes_empty(self) -> None:
        """空字节返回空字符串。"""
        assert TextExtractor().extract_from_bytes(b"") == ""

    def test_text_extract_from_bytes_max_size(self) -> None:
        """超过 max_size 的字节返回空字符串。"""
        extractor = TextExtractor(max_size=10)
        assert extractor.extract_from_bytes(b"x" * 100) == ""

    def test_text_extract_from_bytes_gbk(self, gbk_file: Path) -> None:
        """GBK 编码字节应正确解码。"""
        data = gbk_file.read_bytes()
        content = TextExtractor().extract_from_bytes(data)
        assert "密码" in content
        assert "password123" in content

    def test_docx_extract_from_bytes_matches_path(self, docx_file: Path) -> None:
        """DocxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = docx_file.read_bytes()
        extractor = DocxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(docx_file)

    def test_pptx_extract_from_bytes_matches_path(self, pptx_file: Path) -> None:
        """PptxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = pptx_file.read_bytes()
        extractor = PptxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(pptx_file)

    def test_xlsx_extract_from_bytes_matches_path(self, xlsx_file: Path) -> None:
        """XlsxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = xlsx_file.read_bytes()
        extractor = XlsxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(xlsx_file)

    def test_pdf_extract_from_bytes_matches_path(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """PdfExtractor 从 bytes 提取与从 path 提取结果一致。"""
        # iter-91：强制走 pypdf 回退路径
        monkeypatch.setattr("fuscan.extractors.pdf._PDF_OXIDE_AVAILABLE", False)
        path = tmp_path / "fake.pdf"
        path.write_bytes(b"fake pdf content")

        class FakePage:
            def extract_text(self) -> str:
                return "页面文本 password"

        class FakeReader:
            def __init__(self) -> None:
                self.is_encrypted = False
                self.pages = [FakePage()]

        class FakePdfModule:
            PdfReader = staticmethod(lambda _: FakeReader())

            class errors:
                class PdfReadError(Exception):
                    pass

        import sys

        monkeypatch.setitem(sys.modules, "pypdf", FakePdfModule)
        monkeypatch.setitem(sys.modules, "pypdf.errors", type("errors", (), {"PdfReadError": Exception}))

        data = path.read_bytes()
        extractor = PdfExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_odt_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdtExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from odf.opendocument import OpenDocumentText
        from odf.text import P

        doc = OpenDocumentText()
        doc.text.addElement(P(text="odt password 内容"))  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "test.odt"
        doc.save(str(path))

        data = path.read_bytes()
        extractor = OdtExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_ods_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdsExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableCell, TableRow
        from odf.text import P

        doc = OpenDocumentSpreadsheet()
        table = Table(name="数据")
        row = TableRow()
        cell = TableCell()
        cell.addElement(P(text="ods_password"))
        row.addElement(cell)
        table.addElement(row)
        doc.spreadsheet.addElement(table)  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "test.ods"
        doc.save(str(path))

        data = path.read_bytes()
        extractor = OdsExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_wps_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """WpsExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("wps bytes password")
        path = tmp_path / "test.wps"
        doc.save(str(path))

        data = path.read_bytes()
        extractor = WpsExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_wps_extract_from_bytes_non_zip(self) -> None:
        """非 ZIP 格式（旧版二进制）返回空字符串。"""
        assert WpsExtractor().extract_from_bytes(b"\xd0\xcf\x11\xe0 old binary") == ""


class TestExtractContentFromBytes:
    """extract_content_from_bytes 模块函数测试。"""

    def test_extract_text_from_bytes(self, text_file: Path) -> None:
        """按扩展名从字节提取文本。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, "txt")
        assert "hello password world" in content

    def test_extract_docx_from_bytes(self, docx_file: Path) -> None:
        """按扩展名从字节提取 docx。"""
        data = docx_file.read_bytes()
        content = extract_content_from_bytes(data, "docx")
        assert "段落一 含 password" in content

    def test_extract_unknown_extension_returns_empty(self) -> None:
        """未知扩展名返回空字符串。"""
        assert extract_content_from_bytes(b"content", "xyz") == ""

    def test_extract_extension_case_insensitive(self, text_file: Path) -> None:
        """扩展名大小写不敏感。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, "TXT")
        assert "hello password world" in content

    def test_extract_extension_with_dot(self, text_file: Path) -> None:
        """扩展名带点前缀也能正确处理。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, ".txt")
        assert "hello password world" in content


# ---------------------------------------------------------------------------
# 大文件流式读取
# ---------------------------------------------------------------------------


class TestLargeFileStreaming:
    """TextExtractor 大文件流式读取与编码检测。"""

    def test_large_utf8_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-8 文件应流式解码。"""
        # 构造略大于 10MB 的 UTF-8 文件
        line = "password 行内容 " * 10 + "\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large.txt"
        path.write_text(line * repeat, encoding="utf-8")
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content
        assert content.count("\n") == repeat

    def test_large_gbk_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 GBK 文件应流式解码。"""
        line = "密码 password 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("gbk"))) + 1
        path = tmp_path / "large_gbk.txt"
        path.write_bytes((line * repeat).encode("gbk"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content
        assert "密码" in content

    def test_large_utf8_bom_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-8 BOM 文件应流式解码。"""
        line = "password bom 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large_bom.txt"
        path.write_bytes(b"\xef\xbb\xbf" + (line * repeat).encode("utf-8"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content

    def test_large_utf16_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-16 LE BOM 文件应流式解码。"""
        line = "password utf16 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-16-le"))) + 1
        path = tmp_path / "large_utf16.txt"
        path.write_bytes(b"\xff\xfe" + (line * repeat).encode("utf-16-le"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content

    def test_large_file_crlf_normalized(self, tmp_path: Path) -> None:
        """大文件的 CRLF 行尾应规范化为 LF。"""
        line = "password line\r\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large_crlf.txt"
        path.write_bytes((line * repeat).encode("utf-8"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "\r\n" not in content
        assert "password line\n" in content

    def test_large_file_read_os_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """大文件读取 OSError 时抛出 ExtractorError。"""
        path = tmp_path / "large.txt"
        path.write_bytes(b"x" * (10 * 1024 * 1024 + 1))

        original_open = Path.open

        def mock_open(self: Path, *args, **kwargs):  # type: ignore[no-untyped-def]
            if self == path:
                raise OSError("模拟读取失败")
            return original_open(self, *args, **kwargs)

        monkeypatch.setattr(Path, "open", mock_open)
        with pytest.raises(ExtractorError, match="文件读取失败"):
            TextExtractor().extract(path)

    def test_detect_encoding_utf8_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-8 BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xef\xbb\xbfcontent") == "utf-8-sig"

    def test_detect_encoding_utf16_le_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-16 LE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xff\xfecontent") == "utf-16"

    def test_detect_encoding_utf16_be_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-16 BE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xfe\xffcontent") == "utf-16"

    def test_detect_encoding_utf32_le_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-32 LE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xff\xfe\x00\x00content") == "utf-32"

    def test_detect_encoding_utf32_be_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-32 BE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\x00\x00\xfe\xffcontent") == "utf-32"

    def test_detect_encoding_plain_utf8(self) -> None:
        """_detect_encoding_from_header 对纯 UTF-8（无 BOM）返回 utf-8。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header("纯 UTF-8 内容 password".encode()) == "utf-8"

    def test_detect_encoding_gbk(self) -> None:
        """_detect_encoding_from_header 对 GBK 字节返回 gbk。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header("中文 GBK 内容密码".encode("gbk")) == "gbk"

    def test_detect_encoding_binary_returns_none(self) -> None:
        """_detect_encoding_from_header 对非文本字节返回 None。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        # 0x80 不是任何 BOM 前缀，也不是有效的 UTF-8 起始字节或 GBK 引导字节
        assert _detect_encoding_from_header(b"\x80\x81\x82\x83\x84\x85") is None

    def test_detect_encoding_empty(self) -> None:
        """_detect_encoding_from_header 对空字节返回 utf-8（空字节可被任意编码解码）。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"") == "utf-8"

    def test_large_bytes_decode_skips_charset_normalizer(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """大 bytes（>10MB）用文件头检测编码，跳过 charset-normalizer。"""
        # 构造 >10MB 的 UTF-8 bytes
        data = ("password 内容\n" * 800000).encode("utf-8")
        assert len(data) > 10 * 1024 * 1024

        called = False

        def fake_from_bytes(data: bytes):
            nonlocal called
            called = True
            raise AssertionError("charset-normalizer 不应被调用")

        monkeypatch.setattr("charset_normalizer.from_bytes", fake_from_bytes)
        content = TextExtractor().extract_from_bytes(data)
        assert "password" in content
        assert called is False

    def test_large_bytes_unknown_encoding_fallback_to_charset_normalizer(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """大 bytes 文件头无法确定编码时回退到 charset-normalizer。"""
        # 构造 >10MB 的非 UTF-8/GBK 字节（0x80 不是任何 BOM/UTF-8/GBK 引导字节）
        data = b"\x80\x81\x82\x83" * (3 * 1024 * 1024)
        assert len(data) > 10 * 1024 * 1024

        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(data[:65536]) is None

        # 应回退到 charset-normalizer
        content = TextExtractor().extract_from_bytes(data)
        assert isinstance(content, str)
        assert len(content) > 0


# ---------------------------------------------------------------------------
# RTF 提取器
# ---------------------------------------------------------------------------


@pytest.fixture()
def rtf_file(tmp_path: Path) -> Path:
    """生成包含 password 关键词的 RTF 测试文件。"""
    rtf_content = (
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}"
        r"\f0\fs24 Hello password world\par This is a test.}"
    )
    path = tmp_path / "test.rtf"
    path.write_text(rtf_content, encoding="utf-8")
    return path


class TestRtfExtractor:
    def test_supported_extensions(self) -> None:
        assert RtfExtractor().supported_extensions == ("rtf",)

    def test_extract_rtf_text(self, rtf_file: Path) -> None:
        content = RtfExtractor().extract(rtf_file)
        assert "Hello password world" in content
        assert "This is a test" in content

    def test_extract_from_bytes(self) -> None:
        rtf = r"{\rtf1\ansi Hello password}"
        content = RtfExtractor().extract_from_bytes(rtf.encode("utf-8"))
        assert "Hello password" in content

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            RtfExtractor().extract(tmp_path / "nonexistent.rtf")

    def test_registry_has_rtf_extractor(self) -> None:
        assert isinstance(get_extractor("rtf"), RtfExtractor)

    def test_rtf_import_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """striprtf 未安装时应抛出 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "striprtf.striprtf":
                raise ImportError("No module named 'striprtf'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="striprtf 未安装"):
            RtfExtractor().extract_from_bytes(b"{\\rtf1 fake}")

    def test_rtf_parse_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """rtf_to_text 抛异常时应包装为 ExtractorError。"""

        def raise_parse(text: str) -> None:
            raise RuntimeError("解析失败")

        monkeypatch.setattr("striprtf.striprtf.rtf_to_text", raise_parse)
        with pytest.raises(ExtractorError, match="RTF 解析失败"):
            RtfExtractor().extract_from_bytes(b"{\\rtf1 fake}")


# ---------------------------------------------------------------------------
# EML 提取器
# ---------------------------------------------------------------------------


@pytest.fixture()
def eml_file(tmp_path: Path) -> Path:
    """生成包含 password 关键词的 EML 测试文件。"""
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["Subject"] = "Test Subject"
    msg["From"] = "sender@example.com"
    msg.set_content("Hello password world")
    path = tmp_path / "test.eml"
    path.write_bytes(msg.as_bytes())
    return path


@pytest.fixture()
def html_eml_file(tmp_path: Path) -> Path:
    """生成仅含 HTML 正文的 EML 测试文件。"""
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["Subject"] = "HTML Test"
    msg["From"] = "sender@example.com"
    msg.set_content("<html><body><p>Hello <b>password</b></p></body></html>", subtype="html")
    path = tmp_path / "test_html.eml"
    path.write_bytes(msg.as_bytes())
    return path


class TestEmlExtractor:
    def test_supported_extensions(self) -> None:
        assert EmlExtractor().supported_extensions == ("eml",)

    def test_extract_eml_text(self, eml_file: Path) -> None:
        content = EmlExtractor().extract(eml_file)
        assert "Test Subject" in content
        assert "sender@example.com" in content
        assert "Hello password world" in content

    def test_extract_html_body(self, html_eml_file: Path) -> None:
        content = EmlExtractor().extract(html_eml_file)
        assert "password" in content
        assert "Hello" in content

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """EML 解析失败抛 ExtractorError。"""
        import email

        def raise_parse(*args: object, **kwargs: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(email, "message_from_bytes", raise_parse)
        with pytest.raises(ExtractorError, match="EML 解析失败"):
            EmlExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            EmlExtractor().extract(tmp_path / "nonexistent.eml")

    def test_registry_has_eml_extractor(self) -> None:
        assert isinstance(get_extractor("eml"), EmlExtractor)

    def test_eml_with_attachment_skipped(self, tmp_path: Path) -> None:
        """带附件的 EML 应跳过附件，仅提取正文。"""
        from email.message import EmailMessage

        msg = EmailMessage()
        msg["Subject"] = "带附件"
        msg["From"] = "sender@example.com"
        msg.set_content("正文 password 内容")
        msg.add_attachment(b"attachment data", maintype="application", subtype="octet-stream")
        path = tmp_path / "attach.eml"
        path.write_bytes(msg.as_bytes())

        content = EmlExtractor().extract(path)
        assert "正文 password 内容" in content
        assert "attachment data" not in content

    def test_eml_invalid_charset_plain_fallback(self, tmp_path: Path) -> None:
        """text/plain 的 charset 无效时应回退到 UTF-8 解码。"""
        raw = (
            b"From: a@b.com\r\n"
            b"Subject: charset test\r\n"
            b"Content-Type: text/plain; charset=invalid-charset\r\n"
            b"\r\n"
            b"password text content"
        )
        path = tmp_path / "bad_charset.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "password text content" in content

    def test_eml_invalid_charset_html_fallback(self, tmp_path: Path) -> None:
        """text/html 的 charset 无效时应回退到 UTF-8 解码。"""
        raw = (
            b"From: a@b.com\r\n"
            b"Subject: html charset test\r\n"
            b"Content-Type: text/html; charset=invalid-charset\r\n"
            b"\r\n"
            b"<p>password html</p>"
        )
        path = tmp_path / "bad_html_charset.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "password html" in content

    def test_eml_empty_body_returns_empty(self, tmp_path: Path) -> None:
        """无正文的 EML 仅返回主题和发件人。"""
        raw = b"From: a@b.com\r\nSubject: no body\r\n\r\n"
        path = tmp_path / "no_body.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "no body" in content
        assert "a@b.com" in content

    def test_eml_no_subject_no_sender(self, tmp_path: Path) -> None:
        """无主题和发件人的 EML 仅返回正文。"""
        raw = b"Content-Type: text/plain\r\n\r\nbody password text"
        path = tmp_path / "no_headers.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "body password text" in content
        assert "主题" not in content
        assert "发件人" not in content


# ---------------------------------------------------------------------------
# MSG 提取器
# ---------------------------------------------------------------------------


class TestMsgExtractor:
    def test_supported_extensions(self) -> None:
        assert MsgExtractor().supported_extensions == ("msg",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock extract_msg.Message 验证文本提取逻辑。"""
        import extract_msg

        class FakeMessage:
            subject = "Test Subject"
            sender = "sender@example.com"
            body = "Hello password world"

        monkeypatch.setattr(extract_msg, "Message", lambda data: FakeMessage())
        content = MsgExtractor().extract_from_bytes(b"fake msg data")
        assert "Test Subject" in content
        assert "sender@example.com" in content
        assert "Hello password world" in content

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """MSG 解析失败抛 ExtractorError。"""
        import extract_msg

        def raise_parse(data: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(extract_msg, "Message", raise_parse)
        with pytest.raises(ExtractorError, match="MSG 解析失败"):
            MsgExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="MSG 解析失败"):
            MsgExtractor().extract(tmp_path / "nonexistent.msg")

    def test_registry_has_msg_extractor(self) -> None:
        assert isinstance(get_extractor("msg"), MsgExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取文本。"""
        import extract_msg

        class FakeMessage:
            subject = "路径测试"
            sender = "path@example.com"
            body = "password from path"

        monkeypatch.setattr(extract_msg, "Message", lambda path: FakeMessage())
        path = tmp_path / "test.msg"
        path.write_bytes(b"fake msg")
        content = MsgExtractor().extract(path)
        assert "路径测试" in content
        assert "path@example.com" in content
        assert "password from path" in content

    def test_extract_from_path_parse_error(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 解析失败抛 ExtractorError。"""
        import extract_msg

        def raise_parse(path: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(extract_msg, "Message", raise_parse)
        path = tmp_path / "bad.msg"
        path.write_bytes(b"fake")
        with pytest.raises(ExtractorError, match="MSG 解析失败"):
            MsgExtractor().extract(path)

    def test_msg_import_error_from_path(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 时 extract_msg 未安装应抛 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "extract_msg":
                raise ImportError("No module named 'extract_msg'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        path = tmp_path / "test.msg"
        path.write_bytes(b"fake")
        with pytest.raises(ExtractorError, match="extract-msg 未安装"):
            MsgExtractor().extract(path)

    def test_msg_import_error_from_bytes(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract_from_bytes 时 extract_msg 未安装应抛 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "extract_msg":
                raise ImportError("No module named 'extract_msg'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="extract-msg 未安装"):
            MsgExtractor().extract_from_bytes(b"fake")

    def test_msg_empty_body(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """body 为 None 时仅返回主题和发件人。"""
        import extract_msg

        class FakeMessage:
            subject: str | None = "无正文"
            sender: str | None = "nob@example.com"
            body: str | None = None

        monkeypatch.setattr(extract_msg, "Message", lambda data: FakeMessage())
        content = MsgExtractor().extract_from_bytes(b"fake")
        assert "无正文" in content
        assert "nob@example.com" in content

    def test_msg_no_subject_no_sender(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """subject 和 sender 为 None 时仅返回正文。"""
        import extract_msg

        class FakeMessage:
            subject: str | None = None
            sender: str | None = None
            body: str | None = "password only body"

        monkeypatch.setattr(extract_msg, "Message", lambda data: FakeMessage())
        content = MsgExtractor().extract_from_bytes(b"fake")
        assert "password only body" in content
        assert "主题" not in content
        assert "发件人" not in content


# ---------------------------------------------------------------------------
# XLS 提取器
# ---------------------------------------------------------------------------


class TestXlsExtractor:
    """XLS 提取器测试。

    iter-92 起 XlsExtractor 与 XlsxExtractor 共用 calamine (Rust + PyO3) 后端，
    以下 mock 测试通过 ``monkeypatch`` 替换 ``CalamineWorkbook.from_filelike``
    验证文本提取逻辑；calamine 后端的真实解析由 ``test_extractor_benchmark.py``
    覆盖（XLS 二进制样本难以程序化生成，跳过基准测试）。
    """

    def test_supported_extensions(self) -> None:
        assert XlsExtractor().supported_extensions == ("xls",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock calamine 验证单元格遍历逻辑。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return [["姓名", "密码"], ["张三", "pwd123"]]

        class FakeWorkbook:
            sheet_names = ["Sheet1"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())

        content = XlsExtractor().extract_from_bytes(b"fake xls data")
        assert "姓名" in content
        assert "pwd123" in content
        assert "张三" in content

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """XLS 解析失败抛 ExtractorError。"""
        import python_calamine

        def raise_parse(_filelike: object) -> None:
            raise python_calamine.CalamineError("解析失败")

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", raise_parse)
        with pytest.raises(ExtractorError, match="XLS 解析失败"):
            XlsExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            XlsExtractor().extract(tmp_path / "nonexistent.xls")

    def test_registry_has_xls_extractor(self) -> None:
        assert isinstance(get_extractor("xls"), XlsExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取单元格文本。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return [["user", "password123"]]

        class FakeWorkbook:
            sheet_names = ["Sheet1"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())
        path = tmp_path / "test.xls"
        path.write_bytes(b"fake xls")
        content = XlsExtractor().extract(path)
        assert "password123" in content
        assert "user" in content

    def test_xls_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时应抛出 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            XlsExtractor().extract_from_bytes(b"fake")

    def test_xls_empty_sheet(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """空工作表应返回空字符串。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return []

        class FakeWorkbook:
            sheet_names = ["Empty"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())
        assert XlsExtractor().extract_from_bytes(b"fake") == ""


# ---------------------------------------------------------------------------
# DOC/PPT 提取器
# ---------------------------------------------------------------------------


class TestExtractUtf16leText:
    """测试 _extract_utf16le_text 辅助函数。"""

    def test_extract_ascii_text(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        text = "Hello password world"
        data = text.encode("utf-16-le")
        result = _extract_utf16le_text(data)
        assert "Hello password world" in result

    def test_extract_chinese_text(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        text = "密码 password 测试"
        data = text.encode("utf-16-le")
        result = _extract_utf16le_text(data)
        assert "密码" in result
        assert "password" in result
        assert "测试" in result

    def test_empty_data(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        assert _extract_utf16le_text(b"") == ""
        assert _extract_utf16le_text(b"\x00") == ""

    def test_skip_short_fragments(self) -> None:
        """长度 < 2 的文本片段被过滤。"""
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        # 单字符 A 后跟非文本字节
        data = b"A\x00\x00\x00B\x00\x00\x00"
        result = _extract_utf16le_text(data)
        assert result == ""


class TestDocExtractor:
    def test_supported_extensions(self) -> None:
        assert DocExtractor().supported_extensions == ("doc",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock olefile.OleFileIO 验证 WordDocument 流文本提取。"""
        import olefile

        text = "Hello password world"
        encoded = text.encode("utf-16-le")

        class FakeStream:
            def read(self) -> bytes:
                return encoded

        class FakeOle:
            def exists(self, name: str) -> bool:
                return name == "WordDocument"

            def openstream(self, name: str) -> FakeStream:
                return FakeStream()

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        content = DocExtractor().extract_from_bytes(b"fake doc data")
        assert "Hello password world" in content

    def test_extract_no_worddocument_stream(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OLE 文件无 WordDocument 流时返回空字符串。"""
        import olefile

        class FakeOle:
            def exists(self, name: str) -> bool:
                return False

            def openstream(self, name: str) -> None:
                raise AssertionError("不应调用 openstream")

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        assert DocExtractor().extract_from_bytes(b"fake") == ""

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """DOC 解析失败抛 ExtractorError。"""
        import olefile

        def raise_parse(data: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(olefile, "OleFileIO", raise_parse)
        with pytest.raises(ExtractorError, match="DOC 解析失败"):
            DocExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            DocExtractor().extract(tmp_path / "nonexistent.doc")

    def test_registry_has_doc_extractor(self) -> None:
        assert isinstance(get_extractor("doc"), DocExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取 WordDocument 流文本。"""
        import olefile

        text = "doc password text"
        encoded = text.encode("utf-16-le")

        class FakeStream:
            def read(self) -> bytes:
                return encoded

        class FakeOle:
            def exists(self, name: str) -> bool:
                return name == "WordDocument"

            def openstream(self, name: str) -> FakeStream:
                return FakeStream()

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        path = tmp_path / "test.doc"
        path.write_bytes(b"fake doc")
        content = DocExtractor().extract(path)
        assert "doc password text" in content

    def test_doc_import_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 未安装时应抛出 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "olefile":
                raise ImportError("No module named 'olefile'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="olefile 未安装"):
            DocExtractor().extract_from_bytes(b"fake")


class TestPptExtractor:
    def test_supported_extensions(self) -> None:
        assert PptExtractor().supported_extensions == ("ppt",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock olefile.OleFileIO 验证 PowerPoint Document 流文本提取。"""
        import olefile

        text = "Slide password content"
        encoded = text.encode("utf-16-le")

        class FakeStream:
            def read(self) -> bytes:
                return encoded

        class FakeOle:
            def exists(self, name: str) -> bool:
                return name == "PowerPoint Document"

            def openstream(self, name: str) -> FakeStream:
                return FakeStream()

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        content = PptExtractor().extract_from_bytes(b"fake ppt data")
        assert "Slide password content" in content

    def test_extract_no_powerpoint_stream(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OLE 文件无 PowerPoint Document 流时返回空字符串。"""
        import olefile

        class FakeOle:
            def exists(self, name: str) -> bool:
                return False

            def openstream(self, name: str) -> None:
                raise AssertionError("不应调用 openstream")

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        assert PptExtractor().extract_from_bytes(b"fake") == ""

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """PPT 解析失败抛 ExtractorError。"""
        import olefile

        def raise_parse(data: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(olefile, "OleFileIO", raise_parse)
        with pytest.raises(ExtractorError, match="PPT 解析失败"):
            PptExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            PptExtractor().extract(tmp_path / "nonexistent.ppt")

    def test_registry_has_ppt_extractor(self) -> None:
        assert isinstance(get_extractor("ppt"), PptExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取 PowerPoint Document 流文本。"""
        import olefile

        text = "ppt password slide"
        encoded = text.encode("utf-16-le")

        class FakeStream:
            def read(self) -> bytes:
                return encoded

        class FakeOle:
            def exists(self, name: str) -> bool:
                return name == "PowerPoint Document"

            def openstream(self, name: str) -> FakeStream:
                return FakeStream()

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        path = tmp_path / "test.ppt"
        path.write_bytes(b"fake ppt")
        content = PptExtractor().extract(path)
        assert "ppt password slide" in content

    def test_ppt_import_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 未安装时应抛出 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "olefile":
                raise ImportError("No module named 'olefile'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="olefile 未安装"):
            PptExtractor().extract_from_bytes(b"fake")


# ---------------------------------------------------------------------------
# 新格式集成测试
# ---------------------------------------------------------------------------


class TestScannerWithNewFormats:
    def test_scan_rtf_content(self, rtf_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.WARNING,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="password"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(rtf_file)
        assert result.has_hit

    def test_scan_eml_content(self, eml_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.WARNING,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="password"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(eml_file)
        assert result.has_hit


class TestContentCache:
    """内容提取缓存测试（需求2：避免重复提取导致卡滞）。"""

    def setup_method(self) -> None:
        """每个测试前清空缓存，确保隔离。"""
        clear_content_cache()

    def test_cached_returns_same_content(self, tmp_path: Path) -> None:
        """缓存提取应返回与直接提取相同的内容。"""
        path = tmp_path / "test.txt"
        path.write_text("hello world\npassword=secret\n", encoding="utf-8")

        direct = extract_content_with_fallback(path)
        cached = extract_content_cached(path)
        assert cached == direct
        assert "password=secret" in cached

    def test_second_call_uses_cache(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """第二次调用相同文件不应重复提取。"""
        path = tmp_path / "test.txt"
        path.write_text("content v1\n", encoding="utf-8")

        call_count = {"n": 0}
        original = extract_content_with_fallback

        def counting_fallback(p: Path) -> str:
            call_count["n"] += 1
            return original(p)

        # 模拟 extract_content_cached 内部调用的 extract_content_with_fallback
        monkeypatch.setattr("fuscan.extractors.cache.extract_content_with_fallback", counting_fallback)

        extract_content_cached(path)
        assert call_count["n"] == 1
        # 第二次调用应命中缓存，不触发提取
        extract_content_cached(path)
        assert call_count["n"] == 1

    def test_file_modified_invalidates_cache(self, tmp_path: Path) -> None:
        """文件修改后（mtime/size 变化）缓存应失效。"""
        path = tmp_path / "test.txt"
        path.write_text("v1\n", encoding="utf-8")

        content1 = extract_content_cached(path)
        assert content1 == "v1\n"

        # 修改文件内容（mtime 和 size 都会变化）
        path.write_text("v2 longer content\n", encoding="utf-8")

        content2 = extract_content_cached(path)
        assert content2 == "v2 longer content\n"

    def test_clear_cache_empties_entries(self, tmp_path: Path) -> None:
        """clear_content_cache 应清空所有缓存项。"""
        path = tmp_path / "test.txt"
        path.write_text("cached\n", encoding="utf-8")
        extract_content_cached(path)

        from fuscan.extractors.cache import _CONTENT_CACHE

        assert len(_CONTENT_CACHE) > 0
        clear_content_cache()
        assert len(_CONTENT_CACHE) == 0

    def test_different_files_cached_separately(self, tmp_path: Path) -> None:
        """不同文件应分别缓存。"""
        p1 = tmp_path / "a.txt"
        p1.write_text("aaa\n", encoding="utf-8")
        p2 = tmp_path / "b.txt"
        p2.write_text("bbb\n", encoding="utf-8")

        assert extract_content_cached(p1) == "aaa\n"
        assert extract_content_cached(p2) == "bbb\n"
        # 再次提取应命中各自缓存
        assert extract_content_cached(p1) == "aaa\n"
        assert extract_content_cached(p2) == "bbb\n"

    def test_stat_failure_falls_back_to_uncached(self, tmp_path: Path) -> None:
        """stat 失败时应回退到无缓存提取。"""
        path = tmp_path / "test.txt"
        path.write_text("fallback\n", encoding="utf-8")

        # 正常提取一次填充缓存
        extract_content_cached(path)

        # 删除文件后再次提取，stat 失败应回退到 extract_content_with_fallback
        # extract_content_with_fallback 内部会 try extract_content 失败后 read_text
        # 文件不存在时 read_text 抛 OSError
        path.unlink()
        with pytest.raises(OSError):
            extract_content_cached(path)

    def test_lru_eviction_when_exceeding_max(self, tmp_path: Path) -> None:
        """超过最大缓存数时淘汰最久未使用的项。"""
        from fuscan.extractors.cache import _CONTENT_CACHE, _CONTENT_CACHE_MAX

        for i in range(_CONTENT_CACHE_MAX + 2):
            p = tmp_path / f"f{i}.txt"
            p.write_text(f"content{i}\n", encoding="utf-8")
            extract_content_cached(p)

        # 缓存数不应超过上限
        assert len(_CONTENT_CACHE) <= _CONTENT_CACHE_MAX
